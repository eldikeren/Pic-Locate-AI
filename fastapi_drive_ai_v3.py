from fastapi import FastAPI, UploadFile, File, Form, HTTPException, Request
from fastapi.responses import JSONResponse, StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import List
import time
import asyncio
import os
import uuid
import openai
from google.oauth2.credentials import Credentials
from google_auth_oauthlib.flow import InstalledAppFlow
from googleapiclient.discovery import build
from googleapiclient.errors import HttpError
from PIL import Image, ImageDraw, ImageFont
import io
import torch
import cv2
import numpy as np
from collections import Counter
import pytesseract
import ssl
import urllib3
from dotenv import load_dotenv
from supabase import create_client, Client

# Load environment variables from .env file
load_dotenv()

# Supabase configuration
SUPABASE_URL = "https://gezmablgrepoaamtizts.supabase.co"
SUPABASE_KEY = "eyJhbGciOiJIUzI1NiIsInR5cCI6IkpXVCJ9.eyJpc3MiOiJzdXBhYmFzZSIsInJlZiI6Imdlem1hYmxncmVwb2FhbXRpenRzIiwicm9sZSI6ImFub24iLCJpYXQiOjE3NTkxNzg2MzMsImV4cCI6MjA3NDc1NDYzM30.lJjaubEzeET8OwcHWJ_x_pOAXd8Bc1yDbpdvKianLM0"

# Initialize Supabase client
supabase: Client = create_client(SUPABASE_URL, SUPABASE_KEY)

# Configure SSL to handle Google Drive SSL issues
ssl_context = ssl.create_default_context()
ssl_context.check_hostname = False
ssl_context.verify_mode = ssl.CERT_NONE

# Disable SSL warnings
urllib3.disable_warnings(urllib3.exceptions.InsecureRequestWarning)

# YOLOv8: ultralytics
from ultralytics import YOLO

# CLIP via transformers
from transformers import CLIPProcessor, CLIPModel

# PDF generation
from reportlab.lib.pagesizes import letter, A4
from reportlab.platypus import SimpleDocTemplate, Image as RLImage, Spacer, Paragraph
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_RIGHT
from reportlab.lib.units import inch
from docx import Document
from pptx import Presentation
from pptx.util import Inches

# Room type detection (prioritize locations over people)
ROOM_KEYWORDS = {
    'kitchen': ['kitchen', 'stove', 'oven', 'refrigerator', 'sink', 'cabinet', 'counter', 'microwave', 'dishwasher', 'מטבח', 'כיריים', 'תנור', 'כיור'],
    'bedroom': ['bed', 'bedroom', 'mattress', 'pillow', 'nightstand', 'dresser', 'wardrobe', 'closet', 'חדר שינה', 'מיטה'],
    'living room': ['sofa', 'couch', 'tv', 'television', 'coffee table', 'living room', 'lounge', 'armchair', 'סלון', 'ספה'],
    'bathroom': ['bathroom', 'toilet', 'sink', 'shower', 'bathtub', 'mirror', 'towel', 'שירותים', 'אמבטיה'],
    'dining room': ['dining table', 'chair', 'dining room', 'table', 'dining', 'פינת אוכל', 'חדר אוכל', 'שולחן'],
    'office': ['desk', 'computer', 'office', 'chair', 'monitor', 'keyboard', 'laptop', 'משרד', 'חדר עבודה'],
    'nursery': ['nursery', 'baby', 'child', 'kids', 'crib', 'toy', 'חדר ילדים', 'חדר תינוק'],
    'garden': ['garden', 'plant', 'tree', 'flower', 'outdoor', 'patio', 'lawn', 'גן', 'גינה'],
    'garage': ['garage', 'car', 'vehicle', 'tool', 'workshop', 'חניה', 'מוסך'],
    'balcony': ['balcony', 'terrace', 'מרפסת'],
    'rooftop': ['rooftop', 'roof', 'גג', 'גג עירוני']
}

# ---------------------------
# App & Globals
# ---------------------------
app = FastAPI(title="Google Drive AI Search v3 (with YOLOv8)")

# Add CORS middleware
app.add_middleware(
    CORSMiddleware,
    allow_origins=["http://localhost:4000", "http://127.0.0.1:4000"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

# Global exception handler to prevent crashes
@app.exception_handler(Exception)
async def global_exception_handler(request, exc):
    print(f"❌ Global exception caught: {exc}")
    import traceback
    print(f"📋 Full traceback: {traceback.format_exc()}")
    return JSONResponse(
        status_code=500,
        content={"error": "Internal server error", "detail": str(exc)}
    )

# Add middleware to handle CORS and other issues
@app.middleware("http")
async def add_cors_and_error_handling(request: Request, call_next):
    """Add CORS headers and handle errors gracefully"""
    try:
        response = await call_next(request)
        response.headers["Access-Control-Allow-Origin"] = "*"
        response.headers["Access-Control-Allow-Methods"] = "GET, POST, PUT, DELETE, OPTIONS"
        response.headers["Access-Control-Allow-Headers"] = "*"
        return response
    except Exception as e:
        print(f"❌ Middleware error: {e}")
        return JSONResponse(
            status_code=500,
            content={"error": "Request processing error"},
            headers={"Access-Control-Allow-Origin": "*"}
        )

# Session storage for authentication
auth_sessions = {}

# Connection cache to avoid repeated authentication
_connection_cache = {
    "last_auth_time": None,
    "auth_duration": None,
    "cached_session": None
}

def refresh_credentials_if_needed(creds):
    """Refresh credentials if they are expired or about to expire"""
    if not creds or not creds.expired:
        return creds
    
    try:
        print("🔄 Refreshing expired credentials...")
        from google.auth.transport.requests import Request
        creds.refresh(Request())
        print("✅ Credentials refreshed successfully")
        return creds
    except Exception as e:
        print(f"❌ Failed to refresh credentials: {e}")
        return None

SCOPES = ['https://www.googleapis.com/auth/drive.readonly']

# OpenAI API Configuration
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
if OPENAI_API_KEY:
    openai.api_key = OPENAI_API_KEY
else:
    print("⚠️ Warning: OPENAI_API_KEY environment variable not set. AI features will be disabled.")

def save_credentials_to_session(session_id: str, credentials: Credentials):
    """Save credentials to session storage"""
    auth_sessions[session_id] = {
        'credentials': credentials,
        'timestamp': time.time(),
        'expires_at': time.time() + 3600  # 1 hour
    }

def load_credentials_from_session(session_id: str) -> Credentials:
    """Load credentials from session storage and refresh if needed"""
    if session_id not in auth_sessions:
        return None
    
    session = auth_sessions[session_id]
    creds = session['credentials']
    
    # Try to refresh credentials if they're expired
    refreshed_creds = refresh_credentials_if_needed(creds)
    if refreshed_creds:
        # Update the session with refreshed credentials
        auth_sessions[session_id]['credentials'] = refreshed_creds
        auth_sessions[session_id]['timestamp'] = time.time()
        auth_sessions[session_id]['expires_at'] = time.time() + 3600  # 1 hour
        return refreshed_creds
    else:
        # If refresh failed, clear the session
        del auth_sessions[session_id]
        return None

def clear_session(session_id: str):
    """Clear session data"""
    if session_id in auth_sessions:
        del auth_sessions[session_id]

drive_service = None
image_index = {}  # {file_id: {'name': str, 'embedding': tensor, 'objects': [], 'colors': []}}
collected_images = {}  # Store selected images across searches
search_feedback = {}   # Store search feedback for learning

def create_company_logo():
    """Create the Idan Locations company logo as an image"""
    try:
        # Create logo image with company branding
        img = Image.new('RGB', (400, 120), color='black')
        draw = ImageDraw.Draw(img)
        
        # Try to use a default font, fallback to basic if not available
        try:
            title_font = ImageFont.truetype("arial.ttf", 32)
            subtitle_font = ImageFont.truetype("arial.ttf", 24)
        except:
            title_font = ImageFont.load_default()
            subtitle_font = ImageFont.load_default()
        
        # Company name colors (teal with orange outline)
        teal_color = (0, 128, 128)  # Dark teal
        orange_color = (255, 165, 0)  # Orange
        
        # Draw "Idan" with outline effect
        idan_text = "Idan"
        idan_bbox = draw.textbbox((0, 0), idan_text, font=title_font)
        idan_width = idan_bbox[2] - idan_bbox[0]
        idan_height = idan_bbox[3] - idan_bbox[1]
        
        # Draw outline (orange)
        for dx in [-2, -1, 0, 1, 2]:
            for dy in [-2, -1, 0, 1, 2]:
                if dx != 0 or dy != 0:
                    draw.text((20 + dx, 20 + dy), idan_text, fill=orange_color, font=title_font)
        
        # Draw main text (teal)
        draw.text((20, 20), idan_text, fill=teal_color, font=title_font)
        
        # Draw "Locations" with outline effect
        locations_text = "Locations"
        locations_bbox = draw.textbbox((0, 0), locations_text, font=subtitle_font)
        locations_width = locations_bbox[2] - locations_bbox[0]
        
        # Draw outline (orange)
        for dx in [-1, 0, 1]:
            for dy in [-1, 0, 1]:
                if dx != 0 or dy != 0:
                    draw.text((20 + dx, 60 + dy), locations_text, fill=orange_color, font=subtitle_font)
        
        # Draw main text (teal)
        draw.text((20, 60), locations_text, fill=teal_color, font=subtitle_font)
        
        # Draw a simple location pin icon (orange carrot-like shape)
        pin_x = 300
        pin_y = 40
        
        # Draw pin body (carrot shape)
        pin_points = [
            (pin_x, pin_y + 40),  # Bottom point
            (pin_x - 15, pin_y + 20),  # Left side
            (pin_x - 10, pin_y),  # Top left
            (pin_x + 10, pin_y),  # Top right
            (pin_x + 15, pin_y + 20)  # Right side
        ]
        draw.polygon(pin_points, fill=orange_color)
        
        # Draw leaves on top (green)
        leaf_color = (0, 150, 0)
        draw.ellipse([pin_x - 8, pin_y - 5, pin_x - 2, pin_y + 5], fill=leaf_color)
        draw.ellipse([pin_x + 2, pin_y - 5, pin_x + 8, pin_y + 5], fill=leaf_color)
        
        # Draw camera aperture in center of pin
        aperture_center = (pin_x, pin_y + 15)
        aperture_radius = 8
        draw.ellipse([
            aperture_center[0] - aperture_radius, 
            aperture_center[1] - aperture_radius,
            aperture_center[0] + aperture_radius, 
            aperture_center[1] + aperture_radius
        ], fill='black', outline=orange_color, width=2)
        
        # Save to bytes
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        return img_bytes.getvalue()
        
    except Exception as e:
        print(f"Error creating company logo: {e}")
        # Return a simple text-based logo as fallback
        img = Image.new('RGB', (300, 80), color='black')
        draw = ImageDraw.Draw(img)
        try:
            font = ImageFont.truetype("arial.ttf", 24)
        except:
            font = ImageFont.load_default()
        draw.text((50, 30), "Idan Locations", fill='white', font=font)
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        return img_bytes.getvalue()

def create_placeholder_image():
    """Create a placeholder image when the real image fails to load"""
    try:
        # Create a simple placeholder image
        img = Image.new('RGB', (300, 200), color='lightgray')
        draw = ImageDraw.Draw(img)
        
        # Try to use a default font, fallback to basic if not available
        try:
            font = ImageFont.truetype("arial.ttf", 20)
        except:
            font = ImageFont.load_default()
        
        # Draw placeholder text
        text = "Image\nNot Available"
        bbox = draw.textbbox((0, 0), text, font=font)
        text_width = bbox[2] - bbox[0]
        text_height = bbox[3] - bbox[1]
        
        x = (300 - text_width) // 2
        y = (200 - text_height) // 2
        
        draw.text((x, y), text, fill='darkgray', font=font)
        
        # Convert to bytes
        img_bytes = io.BytesIO()
        img.save(img_bytes, format='PNG')
        img_bytes.seek(0)
        return img_bytes.getvalue()
    except Exception as e:
        print(f"❌ Failed to create placeholder image: {e}")
        # Return a minimal 1x1 pixel image
        return b'\x89PNG\r\n\x1a\n\x00\x00\x00\rIHDR\x00\x00\x00\x01\x00\x00\x00\x01\x08\x02\x00\x00\x00\x90wS\xde\x00\x00\x00\tpHYs\x00\x00\x0b\x13\x00\x00\x0b\x13\x01\x00\x9a\x9c\x18\x00\x00\x00\nIDATx\x9cc```\x00\x00\x00\x04\x00\x01\xdd\x8d\xb4\x1c\x00\x00\x00\x00IEND\xaeB`\x82'

device = "cuda" if torch.cuda.is_available() else "cpu"

# Hebrew-English mapping for enhanced search
HEBREW_ENGLISH_MAPPING = {
    # 1. סוגי חדרים (Room Types)
    'סלון': 'living room',
    'מטבח': 'kitchen', 
    'פינת אוכל': 'dining area',
    'חדר שינה ראשי': 'master bedroom',
    'חדר שינה ילדים': 'children bedroom',
    'חדר עבודה': 'office',
    'חדר רחצה ראשי': 'master bathroom',
    'חדר רחצה אורחים': 'guest bathroom',
    'חדר משחקים': 'playroom',
    'מרפסת סגורה': 'covered balcony',
    'חצר': 'yard',
    'בריכה': 'pool',
    
    # 2. מבנה ותכנון חלל (Structure and Space Planning)
    'סלון פתוח למטבח': 'open living kitchen',
    'מטבח עם אי': 'kitchen with island',
    'תקרה גבוהה': 'high ceiling',
    'תקרת עץ': 'wooden ceiling',
    'חלונות פנורמיים': 'panoramic windows',
    'חדר עם גלריה': 'room with gallery',
    'דלתות הזזה מזכוכית': 'sliding glass doors',
    'מסדרון ארוך': 'long corridor',
    'חדר עם קירות זכוכית': 'room with glass walls',
    'קיר מחיצה דקורטיבי': 'decorative partition wall',
    
    # 3. סוגי ריצוף (Flooring Types)
    'פרקט עץ טבעי': 'natural wood parquet',
    'בטון מוחלק': 'polished concrete',
    'אריחי קרמיקה': 'ceramic tiles',
    'רצפת שיש': 'marble floor',
    'שטיחים מקיר לקיר': 'wall to wall carpet',
    'רצפת טרצו': 'terrazzo floor',
    'רצפת לינוליאום': 'linoleum floor',
    'רצפת עץ למינציה': 'laminate wood floor',
    'רצפת אפוקסי': 'epoxy floor',
    'מרצפות מצוירות': 'painted tiles',
    
    # 4. סגנון עיצובי (Design Style)
    'מודרני': 'modern',
    'כפרי': 'rustic',
    'תעשייתי': 'industrial',
    'סקנדינבי': 'scandinavian',
    'אקלקטי': 'eclectic',
    'מינימליסטי': 'minimalist',
    'רטרו': 'retro',
    'וינטג': 'vintage',
    'בוהו שיק': 'boho chic',
    'עיצוב אורבני': 'urban design',
    
    # 5. חומרי גמר (Finishing Materials)
    'קירות בטון חשוף': 'exposed concrete walls',
    'עץ מלא': 'solid wood',
    'זכוכית מחוסמת': 'tempered glass',
    'פליז מבריק': 'polished brass',
    'ברזל שחור': 'black iron',
    'עור איכותי': 'quality leather',
    'אריחים מאבן טבעית': 'natural stone tiles',
    'טיח מינרלי': 'mineral plaster',
    'טפטים דקורטיביים': 'decorative wallpaper',
    'מתכת אלומיניום': 'aluminum metal',
    
    # 6. תקרות (Ceilings)
    'תקרה עם קורות עץ': 'ceiling with wooden beams',
    'תקרה גבוהה': 'high ceiling',
    'תקרה נמוכה': 'low ceiling',
    'תקרה עם תאורה נסתרת': 'ceiling with hidden lighting',
    'תקרת גבס מעוצבת': 'designed plaster ceiling',
    'תקרה משופעת': 'sloped ceiling',
    'תקרה מקומרת': 'arched ceiling',
    'תקרת זכוכית': 'glass ceiling',
    'תקרה עם חלון גג': 'ceiling with skylight',
    'תקרה אקוסטית': 'acoustic ceiling',
    
    # 7. פתחים ומעברים (Openings and Transitions)
    'חלונות קיר-לקיר': 'wall to wall windows',
    'דלתות זכוכית': 'glass doors',
    'חלונות בלגיים': 'Belgian windows',
    'פתחים עגולים': 'round openings',
    'דלתות עץ מלא': 'solid wood doors',
    'תריסים מעץ': 'wooden shutters',
    'חלונות עם מסגרת מתכת': 'windows with metal frame',
    'חלונות גג': 'skylights',
    'דלתות פלדה': 'steel doors',
    'פתחי תאורה בתקרה': 'ceiling light openings',
    
    # 8. קומות ומבנה (Floors and Structure)
    'דירה בקומת קרקע': 'ground floor apartment',
    'דירה בקומה עליונה': 'top floor apartment',
    'דופלקס': 'duplex',
    'נטהאוז': 'penthouse',
    'לופט': 'loft',
    'וילה עירונית': 'urban villa',
    'קוטג פרברי': 'suburban cottage',
    'בית לשימור': 'heritage house',
    'מבנה תעשייתי': 'industrial building',
    'מבנה מודרני': 'modern building',
    
    # 9. תאורה (Lighting)
    'תאורה טבעית': 'natural lighting',
    'תאורת לד נסתרת': 'hidden LED lighting',
    'מנורות תקרה מעוצבות': 'designed ceiling lights',
    'מנורות קיר': 'wall lights',
    'תאורת אווירה': 'ambient lighting',
    'תאורה מתכווננת': 'adjustable lighting',
    'תאורת פלורסנט': 'fluorescent lighting',
    'נברשות קריסטל': 'crystal chandeliers',
    'תאורת שולחן': 'table lighting',
    'תאורת חוץ': 'outdoor lighting',
    
    # 10. ריהוט (Furniture)
    'ספות עור': 'leather sofas',
    'שולחן אוכל מעץ מלא': 'solid wood dining table',
    'שולחן קפה': 'coffee table',
    'כורסאות בד': 'fabric armchairs',
    'כסאות בר': 'bar stools',
    'מיטה זוגית': 'double bed',
    'ארון קיר': 'wall wardrobe',
    'מדפים פתוחים': 'open shelves',
    'פינת ישיבה': 'seating area',
    'שולחן עבודה': 'work desk',
    
    # 11. קירות (Walls)
    'קיר לבנים חשוף': 'exposed brick wall',
    'קיר עם טפטים': 'wall with wallpaper',
    'קיר צבעוני': 'colored wall',
    'קיר עם אריחים דקורטיביים': 'wall with decorative tiles',
    'קיר עם תמונות אמנות': 'wall with art pictures',
    'קיר עם מראה גדולה': 'wall with large mirror',
    'קיר מעץ טבעי': 'natural wood wall',
    'קיר עם דלתות נסתרות': 'wall with hidden doors',
    'קיר גבס אקוסטי': 'acoustic plaster wall',
    'קיר עם תאורה פנימית': 'wall with internal lighting',
    
    # 12. חצר וגינה (Yard and Garden)
    'גינה פרטית': 'private garden',
    'גינה אורבנית': 'urban garden',
    'מרפסת תלויה': 'hanging balcony',
    'חצר פנימית': 'inner courtyard',
    'גג ירוק': 'green roof',
    'דק עץ': 'wooden deck',
    'אזור ישיבה חיצוני': 'outdoor seating area',
    'בריכת שחייה פרטית': 'private swimming pool',
    'גינה ים-תיכונית': 'Mediterranean garden',
    'שבילי גינה מרוצפים': 'paved garden paths',
    
    # 13. אקססוריז דקורטיביים (Decorative Accessories)
    'כריות נוי': 'decorative cushions',
    'שטיחים מעוצבים': 'designed rugs',
    'תמונות קיר': 'wall pictures',
    'וילונות בד': 'fabric curtains',
    'פסלים מודרניים': 'modern sculptures',
    'נרות ריחניים': 'scented candles',
    'אגרטלים זכוכית': 'glass vases',
    'מראות מעוצבות': 'designed mirrors',
    'כלי חרס': 'ceramic vessels',
    'עציצים וצמחים ירוקים': 'pots and green plants',
    
    # 14. צבעוניות (Color Palette)
    'מונוכרומטי': 'monochromatic',
    'גווני פסטל': 'pastel tones',
    'צבעים ניטרליים': 'neutral colors',
    'פלטת צבעים חמה': 'warm color palette',
    'פלטת צבעים קרה': 'cool color palette',
    'גוונים כהים ועשירים': 'dark and rich tones',
    'קירות צבעוניים': 'colored walls',
    'ניגודי שחור ולבן': 'black and white contrasts',
    'צבעי אדמה': 'earth colors',
    'גוונים בהירים ומאווררים': 'bright and airy tones',
    
    # 15. אופי החלל (Space Character)
    'חלל פתוח': 'open space',
    'חלל רב-תכליתי': 'multi-purpose space',
    'חלל אינטימי': 'intimate space',
    'חלל עם זרימה טבעית': 'space with natural flow',
    'חלל רשמי': 'formal space',
    'חלל מחולק באלגנטיות': 'elegantly divided space',
    'חלל עם נוף פתוח': 'space with open view',
    'חלל מרווח': 'spacious area',
    'חלל קומפקטי': 'compact space',
    'חלל מואר היטב': 'well-lit space',
    
}

# CLIP model via transformers
clip_model = CLIPModel.from_pretrained("openai/clip-vit-base-patch32")
clip_processor = CLIPProcessor.from_pretrained("openai/clip-vit-base-patch32")
clip_model.to(device)

# YOLOv8 model
yolo_model = YOLO("yolov8n.pt")  # small model, replace with custom if needed

# ---------------------------
# Utilities
# ---------------------------
def extract_dominant_colors(image, num_colors=3):
    """Extract dominant colors from image using KMeans clustering"""
    img = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
    img = cv2.resize(img, (100, 100))  # resize for speed
    img = img.reshape((-1, 3)).astype(np.float32)
    
    criteria = (cv2.TERM_CRITERIA_EPS + cv2.TERM_CRITERIA_MAX_ITER, 10, 1.0)
    _, labels, centers = cv2.kmeans(img, num_colors, None, criteria, 10, cv2.KMEANS_RANDOM_CENTERS)
    colors = [tuple(map(int, center)) for center in centers]
    return colors

def detect_objects_yolo(image):
    """Detect objects in image using YOLOv8"""
    results = yolo_model.predict(np.array(image))
    objs = []
    for r in results:
        for box in r.boxes:
            cls = int(box.cls)
            name = yolo_model.names[cls]
            objs.append(name.lower())
    return objs

def extract_text_ocr(image):
    """Extract text from image using OCR"""
    try:
        # Check if Tesseract is available
        try:
            pytesseract.get_tesseract_version()
        except Exception:
            print("⚠️ Tesseract OCR not available, skipping OCR text extraction")
            return ""
        
        # Convert PIL image to OpenCV format
        img_cv = cv2.cvtColor(np.array(image), cv2.COLOR_RGB2BGR)
        
        # Preprocess image for better OCR
        gray = cv2.cvtColor(img_cv, cv2.COLOR_BGR2GRAY)
        
        # Apply threshold to get better text recognition
        _, thresh = cv2.threshold(gray, 0, 255, cv2.THRESH_BINARY + cv2.THRESH_OTSU)
        
        # Extract text using Tesseract
        text = pytesseract.image_to_string(thresh, lang='eng+heb')
        
        # Clean up text
        text = text.strip().replace('\n', ' ').replace('\r', ' ')
        text = ' '.join(text.split())  # Remove extra whitespace
        
        return text if text else ""
    except Exception as e:
        print(f"OCR error: {e}")
        return ""

# Supabase Vector Database Functions
def store_image_embedding(file_id: str, file_name: str, embedding: torch.Tensor, objects: list, colors: list, folder: str, ocr_text: str):
    """Store image embedding and metadata in Supabase"""
    try:
        # Convert tensor to list for JSON storage
        embedding_list = embedding.cpu().numpy().tolist()
        
        # Prepare data for Supabase
        data = {
            "file_id": file_id,
            "file_name": file_name,
            "embedding": embedding_list,
            "objects": objects,
            "colors": colors,
            "folder": folder,
            "ocr_text": ocr_text
        }
        
        # Use upsert to avoid duplicates
        result = supabase.table("image_embeddings").upsert(data, on_conflict="file_id").execute()
        print(f"✅ Stored/Updated embedding for {file_name} in Supabase")
        return True
        
    except Exception as e:
        print(f"❌ Failed to store embedding in Supabase: {e}")
        print(f"   Data: {data}")
        return False

def search_similar_images(query_embedding: torch.Tensor, top_k: int = 10, filters: dict = None):
    """Search for similar images using Supabase vector similarity"""
    try:
        # Convert query embedding to list
        query_vector = query_embedding.cpu().numpy().tolist()
        
        # Build the query
        query = supabase.table("image_embeddings").select("*")
        
        # Add filters if provided
        if filters:
            if filters.get("objects"):
                query = query.contains("objects", filters["objects"])
            if filters.get("folder"):
                query = query.eq("folder", filters["folder"])
        
        # Perform vector similarity search
        # Note: This requires pgvector extension in Supabase
        result = query.execute()
        
        if not result.data:
            return []
        
        # Calculate similarities manually (since pgvector might not be set up)
        similarities = []
        for row in result.data:
            stored_embedding = torch.tensor(row["embedding"])
            similarity = torch.cosine_similarity(query_embedding, stored_embedding, dim=0).item()
            similarities.append((row, similarity))
        
        # Sort by similarity and return top_k
        similarities.sort(key=lambda x: x[1], reverse=True)
        return similarities[:top_k]
        
    except Exception as e:
        print(f"❌ Supabase search failed: {e}")
        return []

def clear_supabase_embeddings():
    """Clear all embeddings from Supabase"""
    try:
        result = supabase.table("image_embeddings").delete().neq("file_id", "").execute()
        print(f"✅ Cleared {len(result.data)} embeddings from Supabase")
        return True
    except Exception as e:
        print(f"❌ Failed to clear Supabase embeddings: {e}")
        return False

def is_image_indexed(file_id: str):
    """Check if an image is already indexed in Supabase"""
    try:
        result = supabase.table("image_embeddings").select("file_id").eq("file_id", file_id).execute()
        return len(result.data) > 0
    except Exception as e:
        print(f"❌ Failed to check if image is indexed: {e}")
        return False

def load_existing_embeddings():
    """Load existing embeddings from Supabase into local index"""
    try:
        result = supabase.table("image_embeddings").select("*").execute()
        if result.data:
            print(f"📥 Loading {len(result.data)} existing embeddings from Supabase...")
            for row in result.data:
                file_id = row["file_id"]
                image_index[file_id] = {
                    "name": row["file_name"],
                    "embedding": torch.tensor(row["embedding"]),
                    "objects": row["objects"],
                    "colors": row["colors"],
                    "folder": row["folder"],
                    "ocr_text": row["ocr_text"]
                }
            print(f"✅ Loaded {len(result.data)} embeddings into local index")
            return len(result.data)
        else:
            print("📭 No existing embeddings found in Supabase")
            return 0
    except Exception as e:
        print(f"❌ Failed to load existing embeddings: {e}")
        return 0

def color_match_score(dominant_colors, target_colors_rgb):
    """Compute color match score between dominant colors and target colors"""
    if not target_colors_rgb or not dominant_colors:
        return 1.0
    
    score = 0
    for target in target_colors_rgb:
        # Calculate distance to closest dominant color
        distances = [float(np.linalg.norm(np.array(c) - np.array(target))) for c in dominant_colors]
        min_distance = min(distances) if distances else 1.0
        
        # Normalize distance (max possible distance in RGB space is sqrt(3*255^2))
        max_distance = np.sqrt(3 * 255**2)
        normalized_distance = min_distance / max_distance
        
        # Convert to similarity score (closer = higher score)
        similarity = max(0, 1 - normalized_distance)
        score += similarity
    
    return score / len(target_colors_rgb)

def detect_room_types(objects):
    """Enhanced room type detection based on detected objects"""
    detected_rooms = []
    object_lower = [obj.lower() for obj in objects]
    object_text = ' '.join(object_lower)
    
    # Enhanced room detection with weighted scoring
    room_scores = {}
    
    # Kitchen indicators (high weight for distinctive items)
    kitchen_score = 0
    kitchen_keywords = ['sink', 'refrigerator', 'oven', 'microwave', 'toaster', 'knife', 'bowl', 'cup', 'bottle', 'wine glass', 'dining table', 'stove', 'dishwasher']
    for keyword in kitchen_keywords:
        if keyword in object_text:
            weight = 3 if keyword in ['sink', 'refrigerator', 'oven', 'stove'] else 1
            kitchen_score += weight
    if kitchen_score > 0:
        room_scores['kitchen'] = kitchen_score
    
    # Bedroom indicators
    bedroom_score = 0
    bedroom_keywords = ['bed', 'pillow', 'lamp', 'clock', 'dresser', 'wardrobe']
    for keyword in bedroom_keywords:
        if keyword in object_text:
            weight = 5 if keyword == 'bed' else 1
            bedroom_score += weight
    if bedroom_score > 0:
        room_scores['bedroom'] = bedroom_score
    
    # Bathroom indicators
    bathroom_score = 0
    bathroom_keywords = ['toilet', 'sink', 'towel', 'soap', 'toothbrush', 'bathtub', 'shower']
    for keyword in bathroom_keywords:
        if keyword in object_text:
            weight = 5 if keyword == 'toilet' else 2 if keyword == 'sink' else 1
            bathroom_score += weight
    if bathroom_score > 0:
        room_scores['bathroom'] = bathroom_score
    
    # Living room indicators
    living_score = 0
    living_keywords = ['couch', 'tv', 'remote', 'coffee table', 'book', 'vase', 'sofa', 'armchair']
    for keyword in living_keywords:
        if keyword in object_text:
            weight = 3 if keyword in ['couch', 'tv', 'sofa'] else 1
            living_score += weight
    if living_score > 0:
        room_scores['living room'] = living_score
    
    # Dining room indicators
    dining_score = 0
    dining_keywords = ['dining table', 'chair', 'wine glass', 'bottle', 'plate', 'fork', 'spoon']
    for keyword in dining_keywords:
        if keyword in object_text:
            weight = 3 if keyword == 'dining table' else 1
            dining_score += weight
    if dining_score > 0:
        room_scores['dining room'] = dining_score
    
    # Office indicators
    office_score = 0
    office_keywords = ['laptop', 'keyboard', 'mouse', 'monitor', 'book', 'chair', 'desk', 'computer']
    for keyword in office_keywords:
        if keyword in object_text:
            weight = 3 if keyword in ['laptop', 'computer', 'desk'] else 1
            office_score += weight
    if office_score > 0:
        room_scores['office'] = office_score
    
    # Return rooms sorted by score (highest first)
    if room_scores:
        sorted_rooms = sorted(room_scores.items(), key=lambda x: x[1], reverse=True)
        detected_rooms = [room for room, score in sorted_rooms if score > 0]
    
    return detected_rooms

def analyze_storyboard_image(img):
    """Analyze a storyboard image and extract features"""
    # CLIP embedding
    inputs = clip_processor(images=img, return_tensors="pt")
    inputs = {k: v.to(device) for k, v in inputs.items()}
    with torch.no_grad():
        image_features = clip_model.get_image_features(**inputs)
        embedding = image_features / image_features.norm(dim=-1, keepdim=True)
    
    # Extract dominant colors
    colors = extract_dominant_colors(img)
    
    # YOLO object detection
    objects = detect_objects_yolo(img)
    
    # Detect room types
    suggested_rooms = detect_room_types(objects)
    
    return {
        'embedding': embedding.cpu(),
        'objects': objects,
        'colors': colors,
        'suggested_rooms': suggested_rooms
    }

def translate_hebrew_query(query):
    """Translate Hebrew terms to English for better CLIP understanding"""
    translated_query = query
    
    # Check for Hebrew terms and translate them
    for hebrew_term, english_term in HEBREW_ENGLISH_MAPPING.items():
        if hebrew_term in query:
            translated_query = translated_query.replace(hebrew_term, english_term)
    
    return translated_query

def calculate_combined_score(semantic_score, object_score, color_score, weights=None):
    """Calculate combined score with configurable weights"""
    if weights is None:
        weights = {"semantic": 0.6, "object": 0.2, "color": 0.2}
    
    # Ensure scores are in valid range [0, 1]
    semantic_score = max(0, min(1, semantic_score))
    object_score = max(0, min(1, object_score))
    color_score = max(0, min(1, color_score))
    
    combined = (
        semantic_score * weights["semantic"] +
        object_score * weights["object"] +
        color_score * weights["color"]
    )
    
    return combined

# ---------------------------
# 1️⃣ Authenticate with Google Drive
# ---------------------------
@app.get("/auth")
def auth_drive():
    """Authenticate with Google Drive using Service Account or OAuth2"""
    global drive_service

    # Check if we already have a working connection
    last_auth_time = _connection_cache.get("last_auth_time")
    if drive_service and last_auth_time:
        import time
        time_since_auth = time.time() - last_auth_time
        if time_since_auth < 300:  # 5 minutes cache
            print("🚀 Using cached Google Drive connection")
            return {
                "status": "authenticated",
                "message": "Using cached Google Drive connection",
                "session_id": _connection_cache.get("cached_session"),
                "method": "cached",
                "cached_duration": int(time_since_auth)
            }

    try:
        # Try service account first (more reliable for server applications)
        service_account_file = "secret-spark-432817-r3-e78bdaac1d51.json"
        if os.path.exists(service_account_file):
            print("🔐 Using service account authentication...")
            from google.oauth2 import service_account
            creds = service_account.Credentials.from_service_account_file(
                service_account_file,
                scopes=['https://www.googleapis.com/auth/drive.readonly']
            )
            drive_service = build('drive', 'v3', credentials=creds)
            
            # Auto-load existing embeddings from Supabase
            print("📥 Auto-loading existing embeddings...")
            loaded_count = load_existing_embeddings()
            if loaded_count > 0:
                print(f"✅ Auto-loaded {loaded_count} existing embeddings")

            # Test the connection with timeout
            try:
                print("🔗 Testing Google Drive connection...")
                import concurrent.futures
                
                # Quick connection test with timeout - check if files are accessible
                def test_connection():
                    try:
                        if drive_service:
                            results = drive_service.files().list(pageSize=1).execute()
                            files = results.get('files', [])
                            print(f"🔍 Service account can access {len(files)} files")
                            # If no files accessible, this is a problem
                            if len(files) == 0:
                                print("⚠️ Service account has no file access - will try OAuth2")
                                return None
                            return results
                        return None
                    except Exception as e:
                        print(f"❌ Connection test failed: {e}")
                        return None
                
                # Use thread pool with timeout for connection test
                with concurrent.futures.ThreadPoolExecutor() as executor:
                    future = executor.submit(test_connection)
                    try:
                        results = future.result(timeout=10)  # 10 second timeout
                        if results is None:
                            raise Exception("Service account has no file access - trying OAuth2")
                    except concurrent.futures.TimeoutError:
                        print("⏰ Connection test timed out, but service account should still work")
                        results = {"files": []}  # Assume connection works
                
                print("✅ Google Drive connection successful")
                session_id = str(uuid.uuid4())
                save_credentials_to_session(session_id, creds)
                
                # Cache the connection
                import time
                _connection_cache["last_auth_time"] = time.time()
                _connection_cache["cached_session"] = session_id

                return {
                    "status": "authenticated",
                    "message": "Successfully connected to Google Drive via Service Account",
                    "session_id": session_id,
                    "method": "service_account"
                }
            except HttpError as e:
                print(f"❌ Service Account HttpError: {e}")
                # Continue to OAuth2 fallback
            except Exception as e:
                print(f"❌ Service Account Exception: {e}")
                # Continue to OAuth2 fallback

        # Try OAuth with NEW client secret first
        new_client_secret_file = "client_secret_132538948811-qim43q7uu42eh2vskk1f4g2n3koa8ong.apps.googleusercontent.com.json"
        if os.path.exists(new_client_secret_file):
            print("🔐 Using NEW OAuth client secret...")
            try:
                print(f"📁 Loading OAuth file: {new_client_secret_file}")
                flow = InstalledAppFlow.from_client_secrets_file(
                    new_client_secret_file,
                    scopes=['https://www.googleapis.com/auth/drive.readonly']
                )
                # Enable offline access for refresh tokens
                flow.redirect_uri = 'http://localhost:8080/callback'
                print("✅ OAuth flow created successfully")
                
                # For web applications, we need to return the authorization URL
                # instead of running a local server
                auth_url, _ = flow.authorization_url(prompt='consent')
                print(f"🌐 OAuth authorization URL: {auth_url}")
                
                return {
                    "status": "oauth_required",
                    "message": "OAuth authorization required",
                    "auth_url": auth_url,
                    "method": "oauth_new",
                    "instructions": "Please visit the auth_url to complete OAuth authentication"
                }
                
                print("🔧 Building Drive service...")
                drive_service = build('drive', 'v3', credentials=creds)
                print("✅ Drive service built")
                
                # Test the connection
                print("🧪 Testing Drive API connection...")
                results = drive_service.files().list(pageSize=1).execute()
                print(f"✅ Drive API test successful - found {len(results.get('files', []))} files")
                
                session_id = str(uuid.uuid4())
                save_credentials_to_session(session_id, creds)
                print(f"💾 Session saved: {session_id}")
                
                # Auto-load existing embeddings from Supabase
                print("📥 Auto-loading existing embeddings...")
                loaded_count = load_existing_embeddings()
                if loaded_count > 0:
                    print(f"✅ Auto-loaded {loaded_count} existing embeddings")

                return {
                    "status": "authenticated",
                    "message": "Successfully connected to Google Drive via NEW OAuth",
                    "session_id": session_id,
                    "method": "oauth_new",
                    "debug": {
                        "oauth_file": new_client_secret_file,
                        "files_found": len(results.get('files', [])),
                        "session_id": session_id
                    }
                }
            except Exception as e:
                print(f"❌ NEW OAuth error: {str(e)}")
                import traceback
                print(f"📋 Full traceback: {traceback.format_exc()}")
                return {
                    "error": f"NEW OAuth authentication failed: {str(e)}",
                    "debug": {
                        "oauth_file": new_client_secret_file,
                        "error_type": type(e).__name__,
                        "traceback": traceback.format_exc()
                    }
                }

        # Fallback to old client secret
        old_client_secret_file = "client_secret_1012576941399-515ln173s773sbrrpn3gtmek0d5vc0u5.apps.googleusercontent.com.json"
        if os.path.exists(old_client_secret_file):
            print("🔐 Using OLD OAuth client secret...")
            try:
                print(f"📁 Loading OAuth file: {old_client_secret_file}")
                flow = InstalledAppFlow.from_client_secrets_file(
                    old_client_secret_file,
                    scopes=['https://www.googleapis.com/auth/drive.readonly']
                )
                # Enable offline access for refresh tokens
                flow.redirect_uri = 'http://localhost:8080/callback'
                print("✅ OAuth flow created successfully")
                
                # For web applications, we need to return the authorization URL
                # instead of running a local server
                auth_url, _ = flow.authorization_url(prompt='consent')
                print(f"🌐 OAuth authorization URL: {auth_url}")
                
                return {
                    "status": "oauth_required",
                    "message": "OAuth authorization required",
                    "auth_url": auth_url,
                    "method": "oauth_old",
                    "instructions": "Please visit the auth_url to complete OAuth authentication"
                }
                
                print("🔧 Building Drive service...")
                drive_service = build('drive', 'v3', credentials=creds)
                print("✅ Drive service built")
                
                # Test the connection
                print("🧪 Testing Drive API connection...")
                results = drive_service.files().list(pageSize=1).execute()
                print(f"✅ Drive API test successful - found {len(results.get('files', []))} files")
                
                session_id = str(uuid.uuid4())
                save_credentials_to_session(session_id, creds)
                print(f"💾 Session saved: {session_id}")
                
                # Auto-load existing embeddings from Supabase
                print("📥 Auto-loading existing embeddings...")
                loaded_count = load_existing_embeddings()
                if loaded_count > 0:
                    print(f"✅ Auto-loaded {loaded_count} existing embeddings")

                return {
                    "status": "authenticated",
                    "message": "Successfully connected to Google Drive via OLD OAuth",
                    "session_id": session_id,
                    "method": "oauth_old",
                    "debug": {
                        "oauth_file": old_client_secret_file,
                        "files_found": len(results.get('files', [])),
                        "session_id": session_id
                    }
                }
            except Exception as e:
                print(f"❌ OLD OAuth error: {str(e)}")
                import traceback
                print(f"📋 Full traceback: {traceback.format_exc()}")
                return {
                    "error": f"OLD OAuth authentication failed: {str(e)}",
                    "debug": {
                        "oauth_file": old_client_secret_file,
                        "error_type": type(e).__name__,
                        "traceback": traceback.format_exc()
                    }
                }

        # Fallback to standard credentials.json
        oauth_file = "credentials.json"
        if os.path.exists(oauth_file):
            print("🔐 Using OAuth with credentials.json...")
            return {
                "error": "OAuth requires browser interaction. Please use the service account or run this in an interactive environment."
            }

    except FileNotFoundError:
        return {"error": "No valid credentials found. Please ensure one of: credentials.json, service account JSON, or client secret JSON is present"}
    except Exception as e:
        return {"error": f"Authentication failed: {str(e)}"}

@app.get("/auth/status")
def check_auth_status(session_id: str | None = None):
    """Check if user is authenticated"""
    global drive_service
    
    if session_id:
        # Check session-based authentication
        creds = load_credentials_from_session(session_id)
        if creds:
            try:
                # Test if credentials are still valid
                test_service = build('drive', 'v3', credentials=creds)
                test_service.files().list(pageSize=1).execute()
                drive_service = test_service  # Update global service
                return {
                    "authenticated": True,
                    "message": "Session is valid and active",
                    "session_id": session_id
                }
            except Exception as e:
                # Credentials expired or invalid
                clear_session(session_id)
                return {
                    "authenticated": False,
                    "message": "Session expired or invalid",
                    "error": str(e)
                }
        else:
            return {
                "authenticated": False,
                "message": "No valid session found"
            }
    else:
        # Check global authentication
        if drive_service:
            try:
                # Test if service is still working
                drive_service.files().list(pageSize=1).execute()
                return {
                    "authenticated": True,
                    "message": "Connected to Google Drive"
                }
            except Exception as e:
                drive_service = None
                return {
                    "authenticated": False,
                    "message": "Connection lost",
                    "error": str(e)
                }
        else:
            return {
                "authenticated": False,
                "message": "Not authenticated"
            }

@app.post("/auth/disconnect")
def disconnect(session_id: str | None = None):
    """Disconnect from Google Drive and clear session"""
    global drive_service
    
    if session_id:
        clear_session(session_id)
    
    drive_service = None
    
    return {
        "status": "disconnected",
        "message": "Successfully disconnected from Google Drive"
    }

# ---------------------------
# 2️⃣ Index Drive Images
# ---------------------------
def crawl_drive_images(service, folder_id='root', folder_path='Root', max_images=999999):
    """Recursively crawl Google Drive and index images"""
    global image_index
    
    # Stop if we've already indexed enough images
    if len(image_index) >= max_images:
        print(f"🛑 Reached limit of {max_images} images, stopping crawl")
        return
    
    print(f"🔍 Crawling folder: {folder_path} (ID: {folder_id})")
    
    # Get images in current folder
    if folder_id == 'root':
        # For root folder, get ALL images in the drive
        query = "mimeType contains 'image/' and trashed=false"
        print(f"   🔍 Searching for ALL images in Google Drive...")
    else:
        # For subfolders, get images in this specific folder
        query = f"'{folder_id}' in parents and mimeType contains 'image/' and trashed=false"
    
    print(f"   🔍 Executing query: {query}")
    try:
        results = service.files().list(q=query, fields="files(id,name,parents)", pageSize=1000).execute()
        files = results.get('files', [])
        print(f"   📸 Found {len(files)} images in {folder_path}")
        
        # Debug: Show first few file names if any found
        if files:
            print(f"   📋 Sample files: {[f['name'] for f in files[:3]]}")
        else:
            print(f"   ⚠️ No files found with query: {query}")
            
    except Exception as e:
        print(f"   ❌ Error executing query: {e}")
        files = []

    for file in files:
        # Stop if we've reached the limit
        if len(image_index) >= max_images:
            print(f"🛑 Reached limit of {max_images} images, stopping processing")
            break
            
        file_id = file['id']
        file_name = file['name']
        
        # Check if image is already indexed
        if is_image_indexed(file_id):
            print(f"   ⏭️ Skipping {file_name} - already indexed")
            continue
        
        # Determine folder path for this image
        if folder_id == 'root':
            # For root query, try to get the actual folder path
            parents = file.get('parents', [])
            if parents:
                # Get the parent folder name
                try:
                    parent_folder = service.files().get(fileId=parents[0], fields="name").execute()
                    folder_path = parent_folder.get('name', 'Unknown Folder')
                except:
                    folder_path = 'Root'
            else:
                folder_path = 'Root'
        else:
            folder_path = folder_path
        
        try:
            # Download image with SSL error handling
            request = service.files().get_media(fileId=file_id)
            try:
                file_content = request.execute()
                file_bytes = io.BytesIO(file_content)
            except Exception as ssl_error:
                if "SSL" in str(ssl_error) or "wrong version number" in str(ssl_error):
                    print(f"   ⚠️ SSL error downloading {file_name}, skipping...")
                    continue
                else:
                    raise ssl_error
            
            try:
                img = Image.open(file_bytes).convert("RGB")
            finally:
                # Ensure the BytesIO object is properly handled
                if hasattr(file_bytes, 'close'):
                    file_bytes.close()
            
            # CLIP embedding
            inputs = clip_processor(images=img, return_tensors="pt")
            inputs = {k: v.to(device) for k, v in inputs.items()}
            with torch.no_grad():
                image_features = clip_model.get_image_features(**inputs)
                embedding = image_features / image_features.norm(dim=-1, keepdim=True)
            
            # Extract dominant colors
            colors = extract_dominant_colors(img)
            
            # YOLO object detection
            objects = detect_objects_yolo(img)
            
            # OCR text extraction
            ocr_text = extract_text_ocr(img)
            
            # Store in local index
            image_index[file_id] = {
                "name": file_name,
                "embedding": embedding.cpu(),
                "objects": objects,
                "colors": colors,
                "folder": folder_path,
                "ocr_text": ocr_text
            }
            
            # Store in Supabase vector database
            store_image_embedding(file_id, file_name, embedding, objects, colors, folder_path, ocr_text)
            
            print(f"Indexed: {file_name} - Objects: {objects} - Colors: {colors}")
            
        except Exception as e:
            error_msg = str(e)
            if "SSL" in error_msg or "wrong version number" in error_msg:
                print(f"   ⚠️ SSL error with {file_name}, skipping...")
            elif "timeout" in error_msg.lower():
                print(f"   ⏰ Timeout downloading {file_name}, skipping...")
            else:
                print(f"   ❌ Failed to process {file_name}: {e}")

    # Crawl subfolders recursively
    query_folders = f"'{folder_id}' in parents and mimeType='application/vnd.google-apps.folder' and trashed=false"
    try:
        folder_results = service.files().list(q=query_folders, fields="files(id,name)").execute()
        subfolders = folder_results.get('files', [])
        
        print(f"   📁 Found {len(subfolders)} subfolders in {folder_path}")
        
        for folder in subfolders:
            print(f"   📂 Subfolder: {folder['name']} (ID: {folder['id']})")
        
        for folder in subfolders:
            new_folder_path = f"{folder_path}/{folder['name']}" if folder_path != 'Root' else folder['name']
            print(f"   🔄 Recursively crawling: {folder['name']}")
            try:
                crawl_drive_images(service, folder_id=folder['id'], folder_path=new_folder_path, max_images=max_images)
            except Exception as e:
                print(f"   ❌ Error crawling folder {folder['name']}: {e}")
                print(f"   ⏭️ Continuing with next folder...")
                continue  # Continue with next folder instead of stopping
    except Exception as e:
        print(f"   ❌ Error getting subfolders: {e}")

@app.post("/index")
def index_drive():
    """Index all images in Google Drive"""
    if not drive_service:
        return {"error": "Not authenticated. Call /auth first."}
    
    global image_index
    image_index = {}  # Reset index
    
    try:
        # Start from root to crawl ALL folders and subfolders
        print(f"🎯 Starting indexing from root folder to crawl ALL folders")
        print(f"🔗 This will index images from all accessible folders in your Google Drive")
        print(f"📊 Indexing ALL images - NO LIMIT!")
        print(f"🔍 This is MANUAL indexing - not automatic on connection")
        crawl_drive_images(drive_service, folder_id='root', folder_path='Root', max_images=999999)
        return {
            "status": "Drive indexed successfully", 
            "total_images": len(image_index),
            "message": f"Indexed {len(image_index)} images with CLIP embeddings, YOLO objects, and color data"
        }
    except Exception as e:
        return {"error": f"Indexing failed: {str(e)}"}

# ---------------------------
# 3️⃣ Search Endpoint
# ---------------------------
class SearchRequest(BaseModel):
    query: str
    required_objects: list = []       # e.g., ["island"]
    required_colors: list = []        # e.g., [[128,0,128]] for purple
    top_k: int = 5
    special_guidelines: str = ""      # Special search guidelines
    feedback_images: list = []        # Images to use as positive feedback
    negative_feedback: list = []      # Images to avoid
    search_session_id: str = ""       # To track search sessions

def search_images_internal(search_request):
    """Internal search function that returns results as list (not JSONResponse)"""
    if not image_index:
        return []
    
    # Translate Hebrew query to English for better CLIP understanding
    translated_query = translate_hebrew_query(search_request["query"])
    
    # Encode text query with CLIP
    text_inputs = clip_processor(text=[translated_query], return_tensors="pt", padding=True)
    text_inputs = {k: v.to(device) for k, v in text_inputs.items()}
    with torch.no_grad():
        text_features = clip_model.get_text_features(**text_inputs)
        text_emb = text_features / text_features.norm(dim=-1, keepdim=True)

    results = []
    for fid, data in image_index.items():
        # Semantic similarity score
        sim = (text_emb.cpu() @ data['embedding'].T).item()
        
        # Object filter score
        if search_request.get("required_objects"):
            obj_score = len(set(data['objects']) & set(search_request["required_objects"])) / max(1, len(search_request["required_objects"]))
        else:
            obj_score = 1.0
        
        # Color filter score
        if search_request.get("required_colors"):
            col_score = color_match_score(data['colors'], search_request["required_colors"])
        else:
            col_score = 1.0
        
        # Combined score using improved algorithm
        final_score = calculate_combined_score(sim, obj_score, col_score)
        
        results.append((fid, data['name'], final_score, data['objects'], data['colors'], sim, obj_score, col_score, data.get('folder', 'Root')))

    # Sort by score and return top results
    results.sort(key=lambda x: x[2], reverse=True)
    
    return [{
        "file_id": r[0],
        "name": r[1],
        "score": round(r[2], 4),
        "objects": r[3],
        "colors": r[4],
        "semantic_score": round(r[5], 4),
        "object_score": round(r[6], 4),
        "color_score": round(r[7], 4),
        "folder": r[8]
    } for r in results[:search_request.get("top_k", 6)]]

@app.post("/search")
def search_images(req: SearchRequest):
    """Search images using semantic similarity, object detection, and color matching"""
    if not image_index:
        return {"error": "No images indexed. Call /index first."}
    
    print(f"🔍 Search request: {req.query}")
    print(f"📊 Total indexed images: {len(image_index)}")
    print(f"📂 Folders with images: {set(img.get('folder', 'Unknown') for img in image_index.values())}")
    
    # Translate Hebrew query to English for better CLIP understanding
    translated_query = translate_hebrew_query(req.query)
    print(f"🔄 Translated query: '{req.query}' -> '{translated_query}'")
    
    # Debug: Show sample of indexed images
    sample_images = list(image_index.items())[:3]
    for fid, data in sample_images:
        print(f"📸 Sample image: {data['name']} - Objects: {data['objects']} - Folder: {data.get('folder', 'Unknown')}")
    
    # Apply special guidelines if provided
    if req.special_guidelines:
        translated_query += f" {req.special_guidelines}"
        print(f"📋 Applied guidelines: {req.special_guidelines}")
    
    # Process feedback images for learning
    feedback_boost = 0
    if req.feedback_images:
        print(f"🎯 Using {len(req.feedback_images)} feedback images for learning")
        # Calculate average embedding of feedback images
        feedback_embeddings = []
        for fid in req.feedback_images:
            if fid in image_index:
                feedback_embeddings.append(image_index[fid]['embedding'])
        
        if feedback_embeddings:
            # Average the feedback embeddings
            avg_feedback = torch.stack(feedback_embeddings).mean(dim=0)
            feedback_boost = 0.3  # Boost for similar images
    
    # Encode text query with CLIP
    text_inputs = clip_processor(text=[translated_query], return_tensors="pt", padding=True)
    text_inputs = {k: v.to(device) for k, v in text_inputs.items()}
    with torch.no_grad():
        text_features = clip_model.get_text_features(**text_inputs)
        text_emb = text_features / text_features.norm(dim=-1, keepdim=True)

    results = []
    for fid, data in image_index.items():
        # Skip negative feedback images
        if fid in req.negative_feedback:
            continue
            
        # Semantic similarity score
        sim = (text_emb.cpu() @ data['embedding'].T).item()
        
        # OCR text matching boost
        ocr_boost = 0
        if data.get('ocr_text') and req.query:
            ocr_text = data['ocr_text'].lower()
            query_lower = req.query.lower()
            
            # Check for exact matches in OCR text
            if query_lower in ocr_text:
                ocr_boost = 0.3  # Significant boost for OCR matches
            
            # Check for partial matches
            query_words = query_lower.split()
            ocr_words = ocr_text.split()
            matching_words = len(set(query_words) & set(ocr_words))
            if matching_words > 0:
                ocr_boost = max(ocr_boost, matching_words * 0.1)
        
        sim = sim + ocr_boost
        
        # Apply feedback boost if available
        if req.feedback_images and feedback_boost > 0:
            feedback_sim = (avg_feedback.cpu() @ data['embedding'].T).item()
            sim = max(sim, sim + feedback_sim * feedback_boost)
        
        # Object filter score - enhanced for room type detection
        if req.required_objects:
            obj_score = len(set(data['objects']) & set(req.required_objects)) / max(1, len(req.required_objects))
        else:
            # Check if query is a room type and boost object detection
            room_boost = 0
            if any(room_term in req.query.lower() for room_term in ['kitchen', 'מטבח', 'bedroom', 'חדר שינה', 'bathroom', 'חדר רחצה']):
                # Boost score for kitchen-related objects when searching for kitchen
                if 'kitchen' in translated_query.lower() or 'מטבח' in req.query:
                    kitchen_objects = ['sink', 'refrigerator', 'oven', 'stove', 'microwave', 'dishwasher', 'knife', 'bowl', 'cup', 'bottle', 'wine glass', 'dining table']
                    kitchen_matches = len(set(data['objects']) & set(kitchen_objects))
                    room_boost = min(0.5, kitchen_matches * 0.1)  # Boost up to 0.5
                # Add similar logic for other room types
            obj_score = 1.0 + room_boost
        
        # Color filter score
        if req.required_colors:
            col_score = color_match_score(data['colors'], req.required_colors)
        else:
            col_score = 1.0
        
        # Combined score using improved algorithm with dynamic weights
        # For room type searches, give more weight to object detection
        if any(room_term in req.query.lower() for room_term in ['kitchen', 'מטבח', 'bedroom', 'חדר שינה', 'bathroom', 'חדר רחצה']):
            weights = {"semantic": 0.4, "object": 0.5, "color": 0.1}  # More weight to objects for room detection
        else:
            weights = {"semantic": 0.6, "object": 0.2, "color": 0.2}  # Default weights
        
        final_score = calculate_combined_score(sim, obj_score, col_score, weights)
        
        results.append((fid, data['name'], final_score, data['objects'], data['colors'], sim, obj_score, col_score, data.get('folder', 'Root')))

    # Sort by score and return top results
    results.sort(key=lambda x: x[2], reverse=True)
    
    # Add folder diversity to prevent bias towards one folder
    if len(results) > req.top_k:
        # Group results by folder
        folder_groups = {}
        for r in results:
            folder = r[8]
            if folder not in folder_groups:
                folder_groups[folder] = []
            folder_groups[folder].append(r)
        
        # Select diverse results from different folders
        diverse_results = []
        max_per_folder = max(1, req.top_k // len(folder_groups)) if folder_groups else req.top_k
        
        for folder, folder_results in folder_groups.items():
            diverse_results.extend(folder_results[:max_per_folder])
        
        # Sort by score again and take top results
        diverse_results.sort(key=lambda x: x[2], reverse=True)
        final_results = diverse_results[:req.top_k]
    else:
        final_results = results
    
    # Debug: Print top results with scores
    print(f"🏆 Top {min(5, len(final_results))} search results:")
    for i, r in enumerate(final_results[:5]):
        print(f"  {i+1}. {r[1]} - Score: {r[2]:.4f} (Semantic: {r[5]:.4f}, Objects: {r[6]:.4f}, Colors: {r[7]:.4f}) - Folder: {r[8]} - Objects: {r[3]}")
    
    return JSONResponse(content=[{
        "file_id": r[0],
        "name": r[1],
        "score": round(r[2], 4),
        "objects": r[3],
        "colors": r[4],
        "semantic_score": round(r[5], 4),
        "object_score": round(r[6], 4),
        "color_score": round(r[7], 4),
        "folder": r[8]
    } for r in final_results])

# ---------------------------
# Advanced Search Features
# ---------------------------

@app.post("/add_to_collection")
async def add_to_collection(request: dict):
    """Add selected images to the collection"""
    global collected_images
    
    file_ids = request.get("file_ids", [])
    file_names = request.get("file_names", [])
    search_session_id = request.get("search_session_id", "default")
    
    if not file_ids:
        return {"error": "No file IDs provided"}
    
    # Initialize collection for this session if not exists
    if search_session_id not in collected_images:
        collected_images[search_session_id] = []
    
    # Add images to collection
    for file_id, file_name in zip(file_ids, file_names):
        if file_id in image_index:
            image_data = image_index[file_id].copy()
            image_data['file_id'] = file_id
            image_data['file_name'] = file_name
            image_data['added_at'] = time.time()
            
            # Avoid duplicates
            if not any(img['file_id'] == file_id for img in collected_images[search_session_id]):
                collected_images[search_session_id].append(image_data)
    
    return {
        "status": "success",
        "message": f"Added {len(file_ids)} images to collection",
        "collection_size": len(collected_images[search_session_id])
    }

@app.get("/get_collection/{session_id}")
async def get_collection(session_id: str):
    """Get collected images for a session"""
    if session_id not in collected_images:
        return {"images": [], "count": 0}
    
    return {
        "images": collected_images[session_id],
        "count": len(collected_images[session_id])
    }

@app.post("/clear_collection")
async def clear_collection(request: dict):
    """Clear collected images for a session"""
    global collected_images
    
    session_id = request.get("session_id", "default")
    if session_id in collected_images:
        del collected_images[session_id]
        return {"status": "success", "message": "Collection cleared"}
    else:
        return {"status": "error", "message": "No collection found for session"}

@app.post("/search_with_feedback")
async def search_with_feedback(request: dict):
    """Perform search with feedback from previous results"""
    query = request.get("query", "")
    guidelines = request.get("guidelines", "")
    feedback_images = request.get("feedback_images", [])
    negative_feedback = request.get("negative_feedback", [])
    search_session_id = request.get("search_session_id", "default")
    top_k = request.get("top_k", 10)
    
    # Create search request with feedback
    search_req = SearchRequest(
        query=query,
        special_guidelines=guidelines,
        feedback_images=feedback_images,
        negative_feedback=negative_feedback,
        search_session_id=search_session_id,
        top_k=top_k
    )
    
    # Store feedback for learning
    if search_session_id not in search_feedback:
        search_feedback[search_session_id] = []
    
    search_feedback[search_session_id].append({
        "query": query,
        "guidelines": guidelines,
        "feedback_images": feedback_images,
        "negative_feedback": negative_feedback,
        "timestamp": time.time()
    })
    
    # Perform search with Supabase
    return await search_images_supabase(search_req)

async def search_images_supabase(req: SearchRequest):
    """Enhanced search using Supabase vector database"""
    print(f"🔍 Supabase search request: {req.query}")
    
    # Translate Hebrew query to English for better CLIP understanding
    translated_query = translate_hebrew_query(req.query)
    print(f"🔄 Translated query: '{req.query}' -> '{translated_query}'")
    
    # Apply special guidelines if provided
    if req.special_guidelines:
        translated_query += f" {req.special_guidelines}"
        print(f"📋 Applied guidelines: {req.special_guidelines}")
    
    # Encode text query with CLIP
    text_inputs = clip_processor(text=[translated_query], return_tensors="pt", padding=True)
    text_inputs = {k: v.to(device) for k, v in text_inputs.items()}
    with torch.no_grad():
        text_features = clip_model.get_text_features(**text_inputs)
        text_emb = text_features / text_features.norm(dim=-1, keepdim=True)
    
    # Try Supabase vector search first
    print("🔍 Attempting Supabase vector search...")
    supabase_results = search_similar_images(text_emb, req.top_k, {
        "objects": req.required_objects if req.required_objects else None,
        "folder": None  # No folder filter for now
    })
    
    if supabase_results:
        print(f"✅ Found {len(supabase_results)} results from Supabase")
        # Convert Supabase results to our format
        results = []
        for row_data, similarity in supabase_results:
            results.append((
                row_data["file_id"],
                row_data["file_name"],
                similarity,
                row_data["objects"],
                row_data["colors"],
                similarity,  # semantic_score
                1.0,  # object_score (already filtered)
                1.0,  # color_score
                row_data["folder"]
            ))
        
        # Apply folder diversity if we have enough results
        if len(results) > req.top_k:
            folder_groups = {}
            for r in results:
                folder = r[8]
                if folder not in folder_groups:
                    folder_groups[folder] = []
                folder_groups[folder].append(r)
            
            diverse_results = []
            max_per_folder = max(1, req.top_k // len(folder_groups)) if folder_groups else req.top_k
            
            for folder, folder_results in folder_groups.items():
                diverse_results.extend(folder_results[:max_per_folder])
            
            diverse_results.sort(key=lambda x: x[2], reverse=True)
            final_results = diverse_results[:req.top_k]
        else:
            final_results = results
        
        print(f"🏆 Top {min(5, len(final_results))} Supabase search results:")
        for i, r in enumerate(final_results[:5]):
            print(f"  {i+1}. {r[1]} - Score: {r[2]:.4f} - Folder: {r[8]} - Objects: {r[3]}")
        
        return JSONResponse(content=[{
            "file_id": r[0],
            "name": r[1],
            "score": round(r[2], 4),
            "objects": r[3],
            "colors": r[4],
            "semantic_score": round(r[5], 4),
            "object_score": round(r[6], 4),
            "color_score": round(r[7], 4),
            "folder": r[8]
        } for r in final_results])
    
    print("⚠️ Supabase search failed, falling back to local search...")
    # Fallback to original search
    return search_images(req)

@app.post("/clear_supabase")
async def clear_supabase():
    """Clear all embeddings from Supabase"""
    try:
        success = clear_supabase_embeddings()
        if success:
            return {"status": "success", "message": "Supabase embeddings cleared successfully"}
        else:
            return {"status": "error", "message": "Failed to clear Supabase embeddings"}
    except Exception as e:
        return {"status": "error", "message": f"Error clearing Supabase: {str(e)}"}

@app.post("/setup_supabase_table")
async def setup_supabase_table():
    """Create the image_embeddings table in Supabase"""
    try:
        # This would typically be done via SQL, but we'll try to create via API
        # For now, just return success - the table should be created manually in Supabase
        return {
            "status": "success", 
            "message": "Please create the 'image_embeddings' table in Supabase with columns: file_id, file_name, embedding, objects, colors, folder, ocr_text, created_at",
            "sql": """
            CREATE TABLE image_embeddings (
                file_id TEXT PRIMARY KEY,
                file_name TEXT,
                embedding VECTOR(512),
                objects TEXT[],
                colors INTEGER[][],
                folder TEXT,
                ocr_text TEXT,
                created_at TIMESTAMP DEFAULT NOW()
            );
            """
        }
    except Exception as e:
        return {"status": "error", "message": f"Error setting up table: {str(e)}"}

@app.get("/test_supabase")
async def test_supabase():
    """Test Supabase connection and check table"""
    try:
        # Test basic connection
        result = supabase.table("image_embeddings").select("count").execute()
        
        # Get table info
        table_info = supabase.table("image_embeddings").select("*").limit(5).execute()
        
        return {
            "status": "success",
            "message": "Supabase connection successful",
            "table_exists": True,
            "sample_data": table_info.data,
            "total_records": len(table_info.data)
        }
    except Exception as e:
        return {
            "status": "error", 
            "message": f"Supabase connection failed: {str(e)}",
            "table_exists": False
        }

@app.post("/export_collection_pdf")
async def export_collection_pdf(request: dict):
    """Export collected images to PDF"""
    session_id = request.get("session_id", "default")
    
    if session_id not in collected_images or not collected_images[session_id]:
        return {"error": "No images in collection"}
    
    # Prepare file data for export
    file_ids = [img['file_id'] for img in collected_images[session_id]]
    file_names = [img['file_name'] for img in collected_images[session_id]]
    
    # Use existing PDF export function
    export_request = {
        "file_ids": file_ids,
        "file_names": file_names
    }
    
    return await export_pdf(export_request)

@app.post("/export_collection_word")
async def export_collection_word(request: dict):
    """Export collected images to Word document"""
    session_id = request.get("session_id", "default")
    
    if session_id not in collected_images or not collected_images[session_id]:
        return {"error": "No images in collection"}
    
    # Prepare file data for export
    file_ids = [img['file_id'] for img in collected_images[session_id]]
    file_names = [img['file_name'] for img in collected_images[session_id]]
    
    # Use existing Word export function
    export_request = {
        "file_ids": file_ids,
        "file_names": file_names,
        "include_proposal": False
    }
    
    return await export_word(export_request)

@app.post("/export_collection_ppt")
async def export_collection_ppt(request: dict):
    """Export collected images to PowerPoint presentation"""
    session_id = request.get("session_id", "default")
    
    if session_id not in collected_images or not collected_images[session_id]:
        return {"error": "No images in collection"}
    
    # Prepare file data for export
    file_ids = [img['file_id'] for img in collected_images[session_id]]
    file_names = [img['file_name'] for img in collected_images[session_id]]
    
    # Use existing PowerPoint export function
    export_request = {
        "file_ids": file_ids,
        "file_names": file_names,
        "include_proposal": False
    }
    
    return await export_ppt(export_request)

# ---------------------------
# 4️⃣ Parse Storyboard / PDF
# ---------------------------
def extract_text_from_pdf(content):
    """Extract text from PDF content"""
    try:
        import PyPDF2
        import io
        
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(content))
        text = ""
        
        for page_num, page in enumerate(pdf_reader.pages):
            try:
                page_text = page.extract_text()
                if page_text:
                    text += page_text + "\n"
            except Exception as e:
                print(f"Error extracting text from page {page_num}: {e}")
                continue
        
        # If no text extracted, try alternative method
        if not text.strip():
            try:
                # Try extracting with different method
                for page in pdf_reader.pages:
                    if hasattr(page, 'get_contents'):
                        content_obj = page.get_contents()
                        if content_obj:
                            text += str(content_obj) + "\n"
            except Exception as e:
                print(f"Alternative PDF extraction failed: {e}")
        
        return text.strip() if text.strip() else "No text could be extracted from PDF"
        
    except ImportError:
        return "PyPDF2 not available for PDF text extraction"
    except Exception as e:
        return f"PDF extraction failed: {str(e)}"

def extract_requirements_from_text(text):
    """Extract design requirements from text using OpenAI API for better Hebrew analysis"""
    try:
        # Preprocess text to improve Hebrew parsing
        # Remove extra whitespace and normalize text
        text = ' '.join(text.split())
        
        # Use OpenAI API for better text analysis (new API format)
        client = openai.OpenAI(api_key=OPENAI_API_KEY)
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {
                    "role": "system",
                    "content": """You are an expert in analyzing Hebrew and English text to extract comprehensive interior design and architectural requirements from storyboards and design documents.

CRITICAL: The text is in Hebrew. Pay special attention to Hebrew words and their meanings. Look for ALL categories mentioned in the text.

Extract the following comprehensive information from the text:

1. ROOM TYPES (סוגי חדרים): סלון, מטבח, פינת אוכל, חדר שינה ראשי, חדר שינה ילדים, חדר עבודה, חדר רחצה ראשי, חדר רחצה אורחים, חדר משחקים, מרפסת סגורה, חצר, בריכה

2. STRUCTURE & SPACE PLANNING (מבנה ותכנון חלל): סלון פתוח למטבח, מטבח עם אי, תקרה גבוהה, תקרת עץ, חלונות פנורמיים, חדר עם גלריה, דלתות הזזה מזכוכית, מסדרון ארוך, חדר עם קירות זכוכית, קיר מחיצה דקורטיבי

3. FLOORING TYPES (סוגי ריצוף): פרקט עץ טבעי, בטון מוחלק, אריחי קרמיקה, רצפת שיש, שטיחים מקיר לקיר, רצפת טרצו, רצפת לינוליאום, רצפת עץ למינציה, רצפת אפוקסי, מרצפות מצוירות

4. DESIGN STYLES (סגנון עיצובי): מודרני, כפרי, תעשייתי, סקנדינבי, אקלקטי, מינימליסטי, רטרו, וינטג', בוהו שיק, עיצוב אורבני

5. FINISHING MATERIALS (חומרי גמר): קירות בטון חשוף, עץ מלא, זכוכית מחוסמת, פליז מבריק, ברזל שחור, עור איכותי, אריחים מאבן טבעית, טיח מינרלי, טפטים דקורטיביים, מתכת אלומיניום

6. CEILINGS (תקרות): תקרה עם קורות עץ, תקרה גבוהה, תקרה נמוכה, תקרה עם תאורה נסתרת, תקרת גבס מעוצבת, תקרה משופעת, תקרה מקומרת, תקרת זכוכית, תקרה עם חלון גג, תקרה אקוסטית

7. OPENINGS & TRANSITIONS (פתחים ומעברים): חלונות קיר-לקיר, דלתות זכוכית, חלונות בלגיים, פתחים עגולים, דלתות עץ מלא, תריסים מעץ, חלונות עם מסגרת מתכת, חלונות גג, דלתות פלדה, פתחי תאורה בתקרה

8. FLOORS & STRUCTURE (קומות ומבנה): דירה בקומת קרקע, דירה בקומה עליונה, דופלקס, פנטהאוז, לופט, וילה עירונית, קוטג' פרברי, בית לשימור, מבנה תעשייתי, מבנה מודרני

9. LIGHTING (תאורה): תאורה טבעית, תאורת לד נסתרת, מנורות תקרה מעוצבות, מנורות קיר, תאורת אווירה, תאורה מתכווננת, תאורת פלורסנט, נברשות קריסטל, תאורת שולחן, תאורת חוץ

10. FURNITURE (ריהוט): ספות עור, שולחן אוכל מעץ מלא, שולחן קפה, כורסאות בד, כסאות בר, מיטה זוגית, ארון קיר, מדפים פתוחים, פינת ישיבה, שולחן עבודה

11. WALLS (קירות): קיר לבנים חשוף, קיר עם טפטים, קיר צבעוני, קיר עם אריחים דקורטיביים, קיר עם תמונות אמנות, קיר עם מראה גדולה, קיר מעץ טבעי, קיר עם דלתות נסתרות, קיר גבס אקוסטי, קיר עם תאורה פנימית

12. YARD & GARDEN (חצר וגינה): גינה פרטית, גינה אורבנית, מרפסת תלויה, חצר פנימית, גג ירוק, דק עץ, אזור ישיבה חיצוני, בריכת שחייה פרטית, גינה ים-תיכונית, שבילי גינה מרוצפים

13. DECORATIVE ACCESSORIES (אקססוריז דקורטיביים): כריות נוי, שטיחים מעוצבים, תמונות קיר, וילונות בד, פסלים מודרניים, נרות ריחניים, אגרטלים זכוכית, מראות מעוצבות, כלי חרס, עציצים וצמחים ירוקים

14. COLOR PALETTE (צבעוניות): מונוכרומטי, גווני פסטל, צבעים ניטרליים, פלטת צבעים חמה, פלטת צבעים קרה, גוונים כהים ועשירים, קירות צבעוניים, ניגודי שחור ולבן, צבעי אדמה, גוונים בהירים ומאווררים

15. SPACE CHARACTER (אופי החלל): חלל פתוח, חלל רב-תכליתי, חלל אינטימי, חלל עם זרימה טבעית, חלל רשמי, חלל מחולק באלגנטיות, חלל עם נוף פתוח, חלל מרווח, חלל קומפקטי, חלל מואר היטב

IMPORTANT: Look for specific locations mentioned like "רחוב מרכזי בעיר", "דירה מודרנית", "חוף ים", "גג עירוני", etc. Be comprehensive and identify ALL relevant categories from the text.

Return ONLY a JSON object with this exact structure:
{
    "location": "primary location or null",
    "style": ["list", "of", "styles"],
    "required_objects": ["list", "of", "objects"],
    "required_colors": ["list", "of", "colors"]
}

If no requirements are found, return empty arrays and null for location."""
                },
                {
                    "role": "user",
                    "content": f"Analyze this text and extract design requirements:\n\n{text}"
                }
            ],
            max_tokens=500,
            temperature=0.1
        )
        
        # Parse the JSON response
        import json
        result = json.loads(response.choices[0].message.content)
        
        # Validate the response structure
        if not isinstance(result, dict):
            raise ValueError("Invalid response format")
        
        # Ensure all required fields exist
        return {
            "location": result.get("location"),
            "style": result.get("style", []),
            "required_objects": result.get("required_objects", []),
            "required_colors": result.get("required_colors", [])
        }
        
    except Exception as e:
        print(f"OpenAI API error: {e}")
        # Fallback to simple keyword matching
        return extract_requirements_fallback(text)

def extract_requirements_fallback(text):
    """Fallback method using comprehensive keyword matching for all 15 categories"""
    text_lower = text.lower()
    
    # 1. Room Types (סוגי חדרים)
    locations = {
        "kitchen": ["kitchen", "cooking", "cook", "stove", "sink", "island", "מטבח", "מטבחון", "כיריים", "תנור", "כיור", "מטבח פתוח"],
        "bedroom": ["bedroom", "bed", "sleep", "master", "guest", "חדר שינה", "חדר שינה ראשי", "חדר שינה ילדים", "מיטה"],
        "living room": ["living room", "lounge", "sitting", "tv", "sofa", "סלון", "חדר מגורים", "סלון מודרני"],
        "dining room": ["dining", "dinner", "table", "eat", "פינת אוכל", "חדר אוכל"],
        "office": ["office", "study", "work", "desk", "computer", "משרד", "חדר עבודה"],
        "bathroom": ["bathroom", "bath", "shower", "toilet", "vanity", "שירותים", "אמבטיה", "חדר רחצה ראשי", "חדר רחצה אורחים"],
        "nursery": ["nursery", "baby", "child", "kids", "crib", "חדר ילדים", "חדר תינוק", "חדר משחקים"],
        "garden": ["garden", "yard", "גן", "גינה", "חצר", "בריכה"],
        "balcony": ["balcony", "terrace", "מרפסת", "מרפסת סגורה"],
        "rooftop": ["rooftop", "roof", "גג", "גג עירוני", "גג עירוני עם נוף"],
        "street": ["street", "רחוב", "רחוב מרכזי", "רחוב אורבני", "רחוב אורבני שוקק"],
        "beach": ["beach", "sea", "חוף", "חוף ים", "ים התיכון", "חוף ים בשעת שקיעה"],
        "city": ["city", "עיר", "תל אביב", "ירושלים", "דירה מודרנית"],
        "apartment": ["apartment", "דירה", "דירה מודרנית", "פנים בית"]
    }
    
    detected_location = None
    for location, keywords in locations.items():
        if any(keyword in text_lower for keyword in keywords):
            detected_location = location
            break
    
    # 2. Structure & Space Planning (מבנה ותכנון חלל)
    structure_keywords = [
        "סלון פתוח למטבח", "מטבח עם אי", "תקרה גבוהה", "תקרת עץ", "חלונות פנורמיים",
        "חדר עם גלריה", "דלתות הזזה מזכוכית", "מסדרון ארוך", "חדר עם קירות זכוכית", "קיר מחיצה דקורטיבי"
    ]
    
    # 3. Flooring Types (סוגי ריצוף)
    flooring_keywords = [
        "פרקט עץ טבעי", "בטון מוחלק", "אריחי קרמיקה", "רצפת שיש", "שטיחים מקיר לקיר",
        "רצפת טרצו", "רצפת לינוליאום", "רצפת עץ למינציה", "רצפת אפוקסי", "מרצפות מצוירות"
    ]
    
    # 4. Design Styles (סגנון עיצובי)
    styles = {
        "modern": ["modern", "contemporary", "sleek", "minimalist", "מודרני", "מודרנית"],
        "rustic": ["rustic", "farmhouse", "country", "wooden", "כפרי", "כפרית"],
        "industrial": ["industrial", "metal", "concrete", "exposed", "תעשייתי", "תעשייתית"],
        "scandinavian": ["scandinavian", "scandi", "nordic", "hygge", "סקנדינבי", "סקנדינבית"],
        "eclectic": ["eclectic", "mixed", "varied", "אקלקטי", "אקלקטית"],
        "minimalist": ["minimalist", "minimal", "clean", "simple", "מינימליסטי", "מינימליסטית"],
        "retro": ["retro", "vintage", "classic", "רטרו", "וינטג"],
        "boho chic": ["bohemian", "boho", "vibrant", "בוהו שיק"],
        "urban design": ["urban", "city", "metropolitan", "עיצוב אורבני"]
    }
    
    # 5. Finishing Materials (חומרי גמר)
    materials_keywords = [
        "קירות בטון חשוף", "עץ מלא", "זכוכית מחוסמת", "פליז מבריק", "ברזל שחור",
        "עור איכותי", "אריחים מאבן טבעית", "טיח מינרלי", "טפטים דקורטיביים", "מתכת אלומיניום"
    ]
    
    # 6. Ceilings (תקרות)
    ceiling_keywords = [
        "תקרה עם קורות עץ", "תקרה גבוהה", "תקרה נמוכה", "תקרה עם תאורה נסתרת", "תקרת גבס מעוצבת",
        "תקרה משופעת", "תקרה מקומרת", "תקרת זכוכית", "תקרה עם חלון גג", "תקרה אקוסטית"
    ]
    
    # 7. Openings & Transitions (פתחים ומעברים)
    openings_keywords = [
        "חלונות קיר-לקיר", "דלתות זכוכית", "חלונות בלגיים", "פתחים עגולים", "דלתות עץ מלא",
        "תריסים מעץ", "חלונות עם מסגרת מתכת", "חלונות גג", "דלתות פלדה", "פתחי תאורה בתקרה"
    ]
    
    # 8. Floors & Structure (קומות ומבנה)
    structure_types = [
        "דירה בקומת קרקע", "דירה בקומה עליונה", "דופלקס", "נטהאוז", "לופט",
        "וילה עירונית", "קוטג' פרברי", "בית לשימור", "מבנה תעשייתי", "מבנה מודרני"
    ]
    
    # 9. Lighting (תאורה)
    lighting_keywords = [
        "תאורה טבעית", "תאורת לד נסתרת", "מנורות תקרה מעוצבות", "מנורות קיר", "תאורת אווירה",
        "תאורה מתכווננת", "תאורת פלורסנט", "נברשות קריסטל", "תאורת שולחן", "תאורת חוץ"
    ]
    
    # 10. Furniture (ריהוט)
    furniture_keywords = [
        "ספות עור", "שולחן אוכל מעץ מלא", "שולחן קפה", "כורסאות בד", "כסאות בר",
        "מיטה זוגית", "ארון קיר", "מדפים פתוחים", "פינת ישיבה", "שולחן עבודה"
    ]
    
    # 11. Walls (קירות)
    wall_keywords = [
        "קיר לבנים חשוף", "קיר עם טפטים", "קיר צבעוני", "קיר עם אריחים דקורטיביים", "קיר עם תמונות אמנות",
        "קיר עם מראה גדולה", "קיר מעץ טבעי", "קיר עם דלתות נסתרות", "קיר גבס אקוסטי", "קיר עם תאורה פנימית"
    ]
    
    # 12. Yard & Garden (חצר וגינה)
    garden_keywords = [
        "גינה פרטית", "גינה אורבנית", "מרפסת תלויה", "חצר פנימית", "גג ירוק",
        "דק עץ", "אזור ישיבה חיצוני", "בריכת שחייה פרטית", "גינה ים-תיכונית", "שבילי גינה מרוצפים"
    ]
    
    # 13. Decorative Accessories (אקססוריז דקורטיביים)
    accessories_keywords = [
        "כריות נוי", "שטיחים מעוצבים", "תמונות קיר", "וילונות בד", "פסלים מודרניים",
        "נרות ריחניים", "אגרטלים זכוכית", "מראות מעוצבות", "כלי חרס", "עציצים וצמחים ירוקים"
    ]
    
    # 14. Color Palette (צבעוניות)
    color_keywords = [
        "מונוכרומטי", "גווני פסטל", "צבעים ניטרליים", "פלטת צבעים חמה", "פלטת צבעים קרה",
        "גוונים כהים ועשירים", "קירות צבעוניים", "ניגודי שחור ולבן", "צבעי אדמה", "גוונים בהירים ומאווררים"
    ]
    
    # 15. Space Character (אופי החלל)
    space_character_keywords = [
        "חלל פתוח", "חלל רב-תכליתי", "חלל אינטימי", "חלל עם זרימה טבעית", "חלל רשמי",
        "חלל מחולק באלגנטיות", "חלל עם נוף פתוח", "חלל מרווח", "חלל קומפקטי", "חלל מואר היטב"
    ]
    
    detected_styles = []
    for style, keywords in styles.items():
        if any(keyword in text_lower for keyword in keywords):
            detected_styles.append(style)
    
    # Detect all other categories
    detected_objects = []
    detected_colors = []
    
    # Check for structure keywords
    for keyword in structure_keywords:
        if keyword in text_lower:
            detected_objects.append(keyword)
    
    # Check for flooring keywords
    for keyword in flooring_keywords:
        if keyword in text_lower:
            detected_objects.append(keyword)
    
    # Check for materials keywords
    for keyword in materials_keywords:
        if keyword in text_lower:
            detected_objects.append(keyword)
    
    # Check for ceiling keywords
    for keyword in ceiling_keywords:
        if keyword in text_lower:
            detected_objects.append(keyword)
    
    # Check for openings keywords
    for keyword in openings_keywords:
        if keyword in text_lower:
            detected_objects.append(keyword)
    
    # Check for structure types
    for keyword in structure_types:
        if keyword in text_lower:
            detected_objects.append(keyword)
    
    # Check for lighting keywords
    for keyword in lighting_keywords:
        if keyword in text_lower:
            detected_objects.append(keyword)
    
    # Check for furniture keywords
    for keyword in furniture_keywords:
        if keyword in text_lower:
            detected_objects.append(keyword)
    
    # Check for wall keywords
    for keyword in wall_keywords:
        if keyword in text_lower:
            detected_objects.append(keyword)
    
    # Check for garden keywords
    for keyword in garden_keywords:
        if keyword in text_lower:
            detected_objects.append(keyword)
    
    # Check for accessories keywords
    for keyword in accessories_keywords:
        if keyword in text_lower:
            detected_objects.append(keyword)
    
    # Check for color keywords
    for keyword in color_keywords:
        if keyword in text_lower:
            detected_colors.append(keyword)
    
    # Check for space character keywords
    for keyword in space_character_keywords:
        if keyword in text_lower:
            detected_objects.append(keyword)
    
    # Add basic object detection for common items
    basic_objects = {
        "island": ["island", "kitchen island", "אי", "אי מטבח"],
        "bed": ["bed", "bedframe", "headboard", "מיטה", "מיטות"],
        "sofa": ["sofa", "couch", "settee", "ספה", "ספות"],
        "table": ["table", "dining table", "coffee table", "שולחן", "שולחנות"],
        "chair": ["chair", "dining chair", "armchair", "כיסא", "כיסאות"],
        "stove": ["stove", "cooktop", "range", "כיריים", "תנור", "תנורים"],
        "sink": ["sink", "faucet", "כיור", "כיורים"],
        "lamp": ["lamp", "lighting", "chandelier", "מנורה", "מנורות", "תאורה"],
        "cabinet": ["cabinet", "cupboard", "storage", "ארון", "ארונות", "ארונות מטבח"],
        "mirror": ["mirror", "reflection", "מראה", "מראות"],
        "rug": ["rug", "carpet", "mat", "שטיח", "שטיחים"],
        "curtain": ["curtain", "drape", "blind", "וילון", "וילונות"],
        "window": ["window", "windows", "חלון", "חלונות", "חלונות גדולים"],
        "door": ["door", "doors", "דלת", "דלתות"],
        "marble": ["marble", "שיש", "אבן שיש"],
        "faucet": ["faucet", "tap", "ברז", "ברזים", "ברז מודרני"]
    }
    
    for obj, keywords in basic_objects.items():
        if any(keyword in text_lower for keyword in keywords):
            detected_objects.append(obj)
    
    # Color detection (Hebrew + English) - Enhanced
    colors = {
        "red": ["red", "crimson", "maroon", "אדום", "אדומה", "אדומים"],
        "blue": ["blue", "navy", "azure", "teal", "כחול", "כחולה", "כחולים"],
        "green": ["green", "emerald", "forest", "mint", "ירוק", "ירוקה", "ירוקים"],
        "yellow": ["yellow", "gold", "amber", "צהוב", "צהובה", "צהובים"],
        "white": ["white", "ivory", "cream", "לבן", "לבנה", "לבנים", "לבנות"],
        "black": ["black", "charcoal", "ebony", "שחור", "שחורה", "שחורים", "שחורות"],
        "gray": ["gray", "grey", "silver", "אפור", "אפורה", "אפורים", "אפורות"],
        "brown": ["brown", "tan", "beige", "חום", "חומה", "חומים", "חומות"],
        "pink": ["pink", "rose", "coral", "ורוד", "ורודה", "ורודים", "ורודות"],
        "purple": ["purple", "violet", "lavender", "סגול", "סגולה", "סגולים", "סגולות"],
        "orange": ["orange", "peach", "apricot", "כתום", "כתומה", "כתומים", "כתומות"],
        "warm colors": ["warm colors", "צבעים חמים", "צבעים חמים"]
    }
    
    detected_colors = []
    for color_name, keywords in colors.items():
        if any(keyword in text_lower for keyword in keywords):
            detected_colors.append(color_name)
    
    # Return comprehensive results
    return {
        "location": detected_location,
        "style": detected_styles,
        "required_objects": detected_objects,
        "required_colors": detected_colors
    }

def generate_ai_proposal(selected_images, requirements=None):
    """Generate AI proposal for selected images"""
    try:
        client = openai.OpenAI(api_key=OPENAI_API_KEY)
        
        # Prepare image descriptions
        image_descriptions = []
        for img in selected_images:
            desc = f"Image: {img['name']} - Objects: {', '.join(img['objects'])} - Colors: {len(img['colors'])} colors - Folder: {img['folder']}"
            image_descriptions.append(desc)
        
        # Create prompt for proposal generation
        prompt = f"""You are an expert interior designer creating a professional proposal for a client. 
        
        Based on the following selected images and requirements, create a comprehensive design proposal:
        
        Selected Images:
        {chr(10).join(image_descriptions)}
        
        Client Requirements:
        {requirements if requirements else "No specific requirements provided"}
        
        Please create a professional proposal that includes:
        1. Executive Summary
        2. Design Concept & Vision
        3. Space Analysis
        4. Recommended Design Elements
        5. Color Palette & Materials
        6. Implementation Timeline
        7. Budget Considerations
        8. Next Steps
        
        Make it professional, detailed, and actionable. Use Hebrew and English as appropriate."""
        
        response = client.chat.completions.create(
            model="gpt-3.5-turbo",
            messages=[
                {"role": "system", "content": "You are a professional interior designer with expertise in creating detailed design proposals."},
                {"role": "user", "content": prompt}
            ],
            max_tokens=2000,
            temperature=0.7
        )
        
        return response.choices[0].message.content
        
    except Exception as e:
        print(f"AI proposal generation failed: {e}")
        return f"""# Design Proposal

## Executive Summary
Based on the selected images, we present a comprehensive design proposal that combines functionality with aesthetic appeal.

## Selected Images Analysis
{chr(10).join([f"- {img['name']}: {', '.join(img['objects'])}" for img in selected_images])}

## Design Concept
The proposed design focuses on creating a cohesive and functional space that meets your requirements.

## Implementation
1. Review selected images
2. Finalize design elements
3. Begin implementation
4. Quality assurance

## Next Steps
Please review this proposal and let us know if you'd like to proceed with any modifications."""

@app.post("/parse_requirements")
async def parse_requirements(file: UploadFile = File(...), guidelines: str = Form(""), feedback_images: str = Form(""), negative_feedback: str = Form("")):
    """Parse storyboard/PDF to extract design requirements with re-search capability"""
    try:
        print(f"📄 Starting PDF parsing for file: {file.filename}")
        
        # Add overall timeout for the entire parsing process
        async def parse_with_timeout():
            return await _parse_requirements_internal(file)
        
        try:
            result = await asyncio.wait_for(parse_with_timeout(), timeout=120.0)  # 2 minute total timeout
            return result
        except asyncio.TimeoutError:
            print("❌ PDF parsing timed out after 2 minutes")
            return JSONResponse(
                status_code=408,
                content={"error": "PDF parsing timed out. Please try with a smaller file or simpler text."}
            )
    except Exception as e:
        print(f"❌ PDF parsing failed: {e}")
        return JSONResponse(
            status_code=500,
            content={"error": f"PDF parsing failed: {str(e)}"}
        )

async def _parse_requirements_internal(file: UploadFile):
    """Internal PDF parsing logic"""
    try:
        # Read file with timeout
        print("📖 Reading file content...")
        content = await asyncio.wait_for(file.read(), timeout=30.0)
        print(f"✅ File read successfully, size: {len(content)} bytes")

        # Determine file type and extract text
        if file.filename.lower().endswith('.pdf'):
            print("📄 Extracting text from PDF...")
            text = extract_text_from_pdf(content)
            print(f"✅ PDF text extracted, length: {len(text)} characters")
        else:
            # Assume text file
            print("📝 Decoding text file...")
            text = content.decode(errors="ignore")
            print(f"✅ Text decoded, length: {len(text)} characters")

        if not text.strip():
            print("❌ No text extracted from file")
            return JSONResponse(
                status_code=400,
                content={"error": "Could not extract text from the uploaded file"}
            )

        # Extract requirements using OpenAI API with timeout
        print("🤖 Starting AI requirements extraction...")
        try:
            parsed = await asyncio.wait_for(
                asyncio.get_event_loop().run_in_executor(None, extract_requirements_from_text, text),
                timeout=60.0
            )
            print(f"✅ AI parsing completed: {parsed}")
        except asyncio.TimeoutError:
            print("❌ AI parsing timed out")
            return JSONResponse(
                status_code=408,
                content={"error": "AI parsing timed out. Please try with a smaller file or simpler text."}
            )
        
        # Auto-search for matching images if requirements were found
        print("🔍 Starting auto-search for matching images...")
        search_results = []
        if parsed and (parsed.get("location") or parsed.get("style") or parsed.get("required_objects") or parsed.get("required_colors")):
            print("✅ Requirements found, proceeding with auto-search")
            try:
                # Create search query from requirements
                search_parts = []
                if parsed and parsed.get("location"):
                    search_parts.append(parsed["location"])
                if parsed and parsed.get("style"):
                    search_parts.extend(parsed["style"])
                if parsed and parsed.get("required_objects"):
                    search_parts.extend(parsed["required_objects"][:3])  # Limit to first 3 objects

                search_query = " ".join(search_parts)
                print(f"🔍 Auto-search query: {search_query}")

                # Perform search if we have images indexed
                if image_index and len(image_index) > 0:
                    print(f"📊 Found {len(image_index)} indexed images, performing search...")
                    # Create search request
                    search_request = {
                        "query": search_query,
                        "required_objects": parsed.get("required_objects", []) if parsed else [],
                        "required_colors": parsed.get("required_colors", []) if parsed else [],
                        "top_k": 6
                    }

                    print(f"🔍 Auto-search objects: {parsed.get('required_objects', []) if parsed else []}")
                    print(f"🔍 Auto-search colors: {parsed.get('required_colors', []) if parsed else []}")

                    # Call search function with timeout
                    try:
                        print("🔍 Executing search...")
                        search_results = await asyncio.wait_for(
                            asyncio.get_event_loop().run_in_executor(None, search_images_internal, search_request),
                            timeout=30.0
                        )
                        print(f"✅ Auto-search found {len(search_results)} results")
                    except asyncio.TimeoutError:
                        print("❌ Auto-search timed out")
                        search_results = []
                else:
                    print("❌ No images indexed for auto-search")
                    
            except Exception as e:
                print(f"❌ Auto-search failed: {e}")
                search_results = []
        else:
            print("❌ No requirements found for auto-search")
        
        print("📤 Preparing response...")
        response_data = {
            "status": "success",
            "filename": file.filename,
            "requirements": parsed,
            "auto_search_results": search_results,
            "search_query": search_query if 'search_query' in locals() else None
        }
        print(f"✅ Response prepared: {len(search_results)} auto-search results")
        
        return JSONResponse(content=response_data)
        
    except Exception as e:
        print(f"❌ Parse requirements error: {e}")
        import traceback
        print(f"📋 Full traceback: {traceback.format_exc()}")
        return JSONResponse(
            status_code=400,
            content={"error": f"Failed to parse file: {str(e)}"}
        )

# ---------------------------
# Additional Utility Endpoints
# ---------------------------
@app.get("/stats")
def get_stats():
    """Get indexing statistics"""
    if not image_index:
        return {"message": "No images indexed yet"}
    
    all_objects = []
    all_colors = []
    for data in image_index.values():
        all_objects.extend(data['objects'])
        all_colors.extend(data['colors'])
    
    object_counts = Counter(all_objects)
    color_counts = Counter([str(c) for c in all_colors])
    
    return {
        "total_images": len(image_index),
        "total_objects_detected": len(object_counts),
        "most_common_objects": dict(object_counts.most_common(10)),
        "most_common_colors": dict(color_counts.most_common(10))
    }

@app.get("/health")
def health_check():
    """Health check endpoint"""
    return {
        "status": "healthy",
        "device": device,
        "authenticated": drive_service is not None,
        "images_indexed": len(image_index)
    }

# ---------------------------
# Health Check
# ---------------------------
@app.get("/image/{file_id}")
async def get_image(file_id: str):
    """Serve image from Google Drive with improved error handling and auto re-authentication"""
    global drive_service
    
    # Check if drive_service is available, if not try to re-authenticate
    if not drive_service:
        print("🔄 Drive service not available, attempting to re-authenticate...")
        try:
            auth_result = auth_drive()
            if auth_result.get("status") == "authenticated":
                print("✅ Successfully re-authenticated with Google Drive")
            else:
                raise HTTPException(status_code=401, detail="Failed to re-authenticate with Google Drive")
        except Exception as e:
            print(f"❌ Re-authentication failed: {e}")
            raise HTTPException(status_code=401, detail="Not authenticated with Google Drive and re-authentication failed")
    
    try:
        print(f"🖼️ Loading image: {file_id}")
        
        # Simple, direct approach with better error handling
        try:
            # Get file metadata with error handling
            try:
                file_metadata = drive_service.files().get(fileId=file_id).execute()
                print(f"✅ File metadata retrieved: {file_metadata.get('name', 'Unknown')}")
            except Exception as meta_error:
                print(f"❌ Failed to get file metadata: {meta_error}")
                # Return placeholder for metadata errors
                placeholder = create_placeholder_image()
                return StreamingResponse(
                    io.BytesIO(placeholder),
                    media_type="image/png",
                    headers={"Content-Disposition": f"inline; filename=metadata_error.png"}
                )
            
            # Download file content with retry logic
            file_content = None
            max_retries = 3
            for attempt in range(max_retries):
                try:
                    print(f"🔄 Download attempt {attempt + 1}/{max_retries} for {file_id}")
                    request = drive_service.files().get_media(fileId=file_id)
                    file_content = request.execute()
                    print(f"✅ Image downloaded successfully, size: {len(file_content)} bytes")
                    break
                except Exception as download_error:
                    print(f"❌ Download attempt {attempt + 1} failed: {download_error}")
                    if attempt == max_retries - 1:
                        print(f"❌ All {max_retries} download attempts failed for {file_id}")
                        # Return placeholder for download errors
                        placeholder = create_placeholder_image()
                        return StreamingResponse(
                            io.BytesIO(placeholder),
                            media_type="image/png",
                            headers={"Content-Disposition": f"inline; filename=download_error.png"}
                        )
                    else:
                        import time
                        time.sleep(1)  # Wait 1 second before retry
            
            # Determine content type
            mime_type = file_metadata.get('mimeType', 'image/jpeg')
            
            return StreamingResponse(
                io.BytesIO(file_content),
                media_type=mime_type,
                headers={"Content-Disposition": f"inline; filename={file_metadata.get('name', 'image')}"}
            )
            
        except Exception as e:
            print(f"❌ Unexpected error in image serving: {e}")
            # Return placeholder image for any error
            try:
                placeholder = create_placeholder_image()
                return StreamingResponse(
                    io.BytesIO(placeholder),
                    media_type="image/png",
                    headers={"Content-Disposition": f"inline; filename=error_placeholder.png"}
                )
            except:
                return JSONResponse(status_code=500, content={"error": "Image unavailable"})
                
    except Exception as e:
        print(f"❌ Unexpected error loading image: {e}")
        return JSONResponse(status_code=500, content={"error": "Internal server error"})

@app.post("/export_pdf")
async def export_pdf(request: dict):
    """Export selected images to PDF"""
    if not drive_service:
        raise HTTPException(status_code=401, detail="Not authenticated with Google Drive")
    
    try:
        file_ids = request.get('file_ids', [])
        file_names = request.get('file_names', [])
        
        if not file_ids:
            raise HTTPException(status_code=400, detail="No file IDs provided")
        
        # Create PDF in memory
        pdf_buffer = io.BytesIO()
        doc = SimpleDocTemplate(pdf_buffer, pagesize=A4)
        story = []
        styles = getSampleStyleSheet()
        
        # Add company logo
        try:
            logo_data = create_company_logo()
            if logo_data and len(logo_data) > 0:
                logo_buffer = io.BytesIO(logo_data)
                logo = RLImage(logo_buffer, width=3*inch, height=0.9*inch)
                story.append(logo)
                story.append(Spacer(1, 20))
                print("✅ Logo added to PDF successfully")
            else:
                print("❌ Logo data is empty, using text fallback")
                title = Paragraph("Idan Locations", styles['Title'])
                story.append(title)
                story.append(Spacer(1, 12))
        except Exception as e:
            print(f"❌ Error adding logo to PDF: {e}")
            import traceback
            traceback.print_exc()
            # Add text title as fallback
            title = Paragraph("Idan Locations", styles['Title'])
            story.append(title)
            story.append(Spacer(1, 12))
        
        # Add Hebrew introduction
        hebrew_intro = """
        הצעת לוקיישנים לצילומים
        
        תקציר מנהלים:
        שמחים להציג בפניכם מבחר לוקיישנים שנבחרו במיוחד בהתאם לדרישות ההפקה שלכם. ההצעה נבנתה מתוך מטרה לאפשר לכם מבט ממוקד, נגיש וברור על מקומות פוטנציאליים לצילומים.
        
        חזון וקונספט:
        ב־Idan Locations אנו מתמחים בהתאמת לוקיישנים מדויקים להפקות קולנוע, טלוויזיה ופרסומות. החזון שלנו הוא לחבר בין צרכי ההפקה שלכם לבין המרחב המתאים ביותר מבחינה ויזואלית, לוגיסטית והפקתית.
        
        סקירת לוקיישנים:
        התמונות שלפניכם מציגות אתרים רלוונטיים שנבחרו בקפידה, מתוך שיקולים של נראות, נגישות ותנאי הפקה. כל מיקום נותן מענה לאופי הסצנות והאווירה שברצונכם ליצור.
        
        יתרונות מרכזיים:
        • מגוון סגנונות ונופים במקום אחד
        • נגישות גבוהה לצוותי צילום והפקה
        • אפשרויות גמישות בהתאם לדרישות ההפקה
        • ניסיון וליווי מקצועי לאורך כל התהליך
        
        שלבים הבאים:
        נשמח לקיים פגישת המשך לבחירת הלוקיישן המתאים ביותר ולהתחלת תהליך התיאום בשטח.
        
        תודה על שיתוף הפעולה,
        Idan Locations
        """
        
        # Create custom style for Hebrew text
        hebrew_style = ParagraphStyle(
            'HebrewStyle',
            parent=styles['Normal'],
            fontSize=11,
            spaceAfter=15,
            alignment=TA_RIGHT,  # Right-to-left for Hebrew
            fontName='Helvetica'
        )
        
        intro_paragraph = Paragraph(hebrew_intro, hebrew_style)
        story.append(intro_paragraph)
        story.append(Spacer(1, 20))
        
        # Add section title
        section_title = Paragraph("תמונות נבחרות", styles['Heading1'])
        story.append(section_title)
        story.append(Spacer(1, 12))
        
        for i, (file_id, file_name) in enumerate(zip(file_ids, file_names)):
            try:
                # Get image from Drive or uploaded images
                if file_id.startswith('uploaded_'):
                    # Handle uploaded image
                    if file_id not in image_index or not image_index[file_id].get("is_uploaded"):
                        raise Exception(f"Uploaded image {file_id} not found")
                    image_content = io.BytesIO(image_index[file_id]["image_data"])
                else:
                    # Handle Google Drive image
                    if not drive_service:
                        raise Exception("Not authenticated with Google Drive")
                    request_drive = drive_service.files().get_media(fileId=file_id)
                    image_content = io.BytesIO()
                    downloader = request_drive.execute()
                    image_content.write(downloader)
                    image_content.seek(0)
                
                # Create PIL image to get dimensions
                pil_image = Image.open(image_content)
                image_content.seek(0)
                
                # Calculate size for PDF (max width 6 inches, maintain aspect ratio)
                max_width = 6 * inch
                width, height = pil_image.size
                aspect_ratio = height / width
                
                if width > max_width:
                    display_width = max_width
                    display_height = display_width * aspect_ratio
                else:
                    display_width = width
                    display_height = height
                
                # Add image name
                story.append(Paragraph(f"<b>{file_name}</b>", styles['Normal']))
                story.append(Spacer(1, 6))
                
                # Add image
                rl_image = RLImage(image_content, width=display_width, height=display_height)
                story.append(rl_image)
                story.append(Spacer(1, 12))
                
            except Exception as e:
                # Add error message for failed images
                story.append(Paragraph(f"<b>{file_name}</b> - Error loading image: {str(e)}", styles['Normal']))
                story.append(Spacer(1, 12))
        
        # Build PDF
        doc.build(story)
        pdf_buffer.seek(0)
        
        return StreamingResponse(
            io.BytesIO(pdf_buffer.getvalue()),
            media_type="application/pdf",
            headers={"Content-Disposition": "attachment; filename=Idan_Locations_Proposal.pdf"}
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"PDF generation failed: {str(e)}")

@app.post("/export_word")
async def export_word(request: dict):
    """Export selected images to Word document with AI proposal"""
    try:
        file_ids = request.get("file_ids", [])
        file_names = request.get("file_names", [])
        include_proposal = request.get("include_proposal", True)
        
        if not file_ids:
            return JSONResponse(
                status_code=400,
                content={"error": "No images selected for export"}
            )
        
        # Create Word document
        doc = Document()
        
        # Add company logo
        try:
            logo_data = create_company_logo()
            logo_buffer = io.BytesIO(logo_data)
            logo_paragraph = doc.add_paragraph()
            logo_run = logo_paragraph.runs[0] if logo_paragraph.runs else logo_paragraph.add_run()
            logo_paragraph.alignment = 1  # Center alignment
            # Note: Adding images to Word requires more complex handling
            # For now, we'll add the company name as a styled heading
        except Exception as e:
            print(f"Error adding logo to Word: {e}")
        
        # Add company title
        title = doc.add_heading('Idan Locations', 0)
        title.alignment = 1  # Center alignment
        
        # Add Hebrew introduction
        hebrew_intro = """
        הצעת לוקיישנים לצילומים
        
        תקציר מנהלים:
        שמחים להציג בפניכם מבחר לוקיישנים שנבחרו במיוחד בהתאם לדרישות ההפקה שלכם. ההצעה נבנתה מתוך מטרה לאפשר לכם מבט ממוקד, נגיש וברור על מקומות פוטנציאליים לצילומים.
        
        חזון וקונספט:
        ב־Idan Locations אנו מתמחים בהתאמת לוקיישנים מדויקים להפקות קולנוע, טלוויזיה ופרסומות. החזון שלנו הוא לחבר בין צרכי ההפקה שלכם לבין המרחב המתאים ביותר מבחינה ויזואלית, לוגיסטית והפקתית.
        
        סקירת לוקיישנים:
        התמונות שלפניכם מציגות אתרים רלוונטיים שנבחרו בקפידה, מתוך שיקולים של נראות, נגישות ותנאי הפקה. כל מיקום נותן מענה לאופי הסצנות והאווירה שברצונכם ליצור.
        
        יתרונות מרכזיים:
        • מגוון סגנונות ונופים במקום אחד
        • נגישות גבוהה לצוותי צילום והפקה
        • אפשרויות גמישות בהתאם לדרישות ההפקה
        • ניסיון וליווי מקצועי לאורך כל התהליך
        
        שלבים הבאים:
        נשמח לקיים פגישת המשך לבחירת הלוקיישן המתאים ביותר ולהתחלת תהליך התיאום בשטח.
        
        תודה על שיתוף הפעולה,
        Idan Locations
        """
        
        intro_paragraph = doc.add_paragraph(hebrew_intro)
        intro_paragraph.alignment = 2  # Right alignment for Hebrew
        doc.add_paragraph()  # Add spacing
        
        # Add images section
        doc.add_heading('תמונות נבחרות', level=1)
        
        # Add images
        for file_id, name in zip(file_ids, file_names):
            try:
                # Get image content
                if file_id.startswith("uploaded_"):
                    # Handle uploaded images
                    image_path = f"uploaded_images/{file_id}.jpg"
                    if os.path.exists(image_path):
                        doc.add_paragraph(f"Image: {name}")
                        doc.add_picture(image_path, width=Inches(4))
                else:
                    # Handle Google Drive images
                    if drive_service:
                        request_drive = drive_service.files().get_media(fileId=file_id)
                        image_content = request_drive.execute()
                        
                        # Save temporary image
                        temp_path = f"temp_{file_id}.jpg"
                        with open(temp_path, 'wb') as f:
                            f.write(image_content)
                        
                        doc.add_paragraph(f"Image: {name}")
                        doc.add_picture(temp_path, width=Inches(4))
                        
                        # Clean up temp file
                        os.remove(temp_path)
            except Exception as e:
                print(f"Error adding image {name}: {e}")
                continue
        
        # Save to buffer
        buffer = io.BytesIO()
        doc.save(buffer)
        buffer.seek(0)
        
        return StreamingResponse(
            io.BytesIO(buffer.read()),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            headers={"Content-Disposition": "attachment; filename=idan_locations_proposal.docx"}
        )
        
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"error": f"Word export failed: {str(e)}"}
        )

@app.post("/export_ppt")
async def export_ppt(request: dict):
    """Export selected images to PowerPoint presentation with AI proposal"""
    try:
        file_ids = request.get("file_ids", [])
        file_names = request.get("file_names", [])
        include_proposal = request.get("include_proposal", True)
        
        if not file_ids:
            return JSONResponse(
                status_code=400,
                content={"error": "No images selected for export"}
            )
        
        # Create PowerPoint presentation
        prs = Presentation()
        
        # Title slide with company branding and logo
        title_slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        # Add company logo to title slide
        try:
            logo_data = create_company_logo()
            if logo_data and len(logo_data) > 0:
                # Save logo temporarily
                temp_logo_path = "temp_logo.png"
                with open(temp_logo_path, 'wb') as f:
                    f.write(logo_data)
                
                # Add logo to slide (positioned at top)
                slide.shapes.add_picture(temp_logo_path, Inches(1), Inches(0.5), Inches(8), Inches(2.4))
                
                # Clean up temp file
                if os.path.exists(temp_logo_path):
                    os.remove(temp_logo_path)
                print("✅ Logo added to PowerPoint successfully")
            else:
                print("❌ Logo data is empty")
        except Exception as e:
            print(f"❌ Error adding logo to PowerPoint: {e}")
            import traceback
            traceback.print_exc()
        
        title.text = "הצעת לוקיישנים לצילומים"
        subtitle.text = "Idan Locations"
        
        # Add executive summary slide
        summary_slide_layout = prs.slide_layouts[1]
        summary_slide = prs.slides.add_slide(summary_slide_layout)
        summary_title = summary_slide.shapes.title
        summary_content = summary_slide.placeholders[1]
        
        summary_title.text = "תקציר מנהלים"
        hebrew_summary = """שמחים להציג בפניכם מבחר לוקיישנים שנבחרו במיוחד בהתאם לדרישות ההפקה שלכם. ההצעה נבנתה מתוך מטרה לאפשר לכם מבט ממוקד, נגיש וברור על מקומות פוטנציאליים לצילומים.

חזון וקונספט:
ב־Idan Locations אנו מתמחים בהתאמת לוקיישנים מדויקים להפקות קולנוע, טלוויזיה ופרסומות. החזון שלנו הוא לחבר בין צרכי ההפקה שלכם לבין המרחב המתאים ביותר מבחינה ויזואלית, לוגיסטית והפקתית."""
        
        summary_content.text = hebrew_summary
        
        # Add locations overview slide
        overview_slide_layout = prs.slide_layouts[1]
        overview_slide = prs.slides.add_slide(overview_slide_layout)
        overview_title = overview_slide.shapes.title
        overview_content = overview_slide.placeholders[1]
        
        overview_title.text = "סקירת לוקיישנים ויתרונות"
        hebrew_overview = """סקירת לוקיישנים:
התמונות שלפניכם מציגות אתרים רלוונטיים שנבחרו בקפידה, מתוך שיקולים של נראות, נגישות ותנאי הפקה.

יתרונות מרכזיים:
• מגוון סגנונות ונופים במקום אחד
• נגישות גבוהה לצוותי צילום והפקה
• אפשרויות גמישות בהתאם לדרישות ההפקה
• ניסיון וליווי מקצועי לאורך כל התהליך

שלבים הבאים:
נשמח לקיים פגישת המשך לבחירת הלוקיישן המתאים ביותר ולהתחלת תהליך התיאום בשטח."""
        
        overview_content.text = hebrew_overview
        
        # No AI proposal slides - using only Hebrew text
        
        # Images slides
        for file_id, name in zip(file_ids, file_names):
            try:
                # Create slide for each image
                img_slide_layout = prs.slide_layouts[5]  # Blank layout
                slide = prs.slides.add_slide(img_slide_layout)
                
                # Add title
                title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
                title_frame = title_box.text_frame
                title_frame.text = name
                
                # Get image content
                if file_id.startswith("uploaded_"):
                    # Handle uploaded images
                    image_path = f"uploaded_images/{file_id}.jpg"
                    if os.path.exists(image_path):
                        slide.shapes.add_picture(image_path, Inches(1), Inches(1.5), Inches(8), Inches(6))
                else:
                    # Handle Google Drive images
                    if drive_service:
                        request_drive = drive_service.files().get_media(fileId=file_id)
                        image_content = request_drive.execute()
                        
                        # Save temporary image
                        temp_path = f"temp_{file_id}.jpg"
                        with open(temp_path, 'wb') as f:
                            f.write(image_content)
                        
                        slide.shapes.add_picture(temp_path, Inches(1), Inches(1.5), Inches(8), Inches(6))
                        
                        # Clean up temp file
                        os.remove(temp_path)
            except Exception as e:
                print(f"Error adding image {name}: {e}")
                continue
        
        # Save to buffer
        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        
        return StreamingResponse(
            io.BytesIO(buffer.read()),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": "attachment; filename=idan_locations_proposal.pptx"}
        )
        
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"error": f"PowerPoint export failed: {str(e)}"}
        )

@app.post("/upload_images")
async def upload_images(images: List[UploadFile] = File(...), guidelines: str = Form(""), feedback_images: str = Form(""), negative_feedback: str = Form("")):
    """Upload and index multiple images with re-search capability"""
    global image_index
    
    if not images:
        raise HTTPException(status_code=400, detail="No images provided")
    
    uploaded_images = []
    
    for image_file in images:
        try:
            # Read image file
            image_content = await image_file.read()
            image_bytes = io.BytesIO(image_content)
            
            # Open and process image
            img = Image.open(image_bytes).convert("RGB")
            
            # Generate unique ID for uploaded image
            import uuid
            file_id = f"uploaded_{uuid.uuid4().hex}"
            
            # CLIP embedding
            inputs = clip_processor(images=img, return_tensors="pt")
            inputs = {k: v.to(device) for k, v in inputs.items()}
            with torch.no_grad():
                image_features = clip_model.get_image_features(**inputs)
                embedding = image_features / image_features.norm(dim=-1, keepdim=True)
            
            # Extract dominant colors
            colors = extract_dominant_colors(img)
            
            # YOLO object detection
            objects = detect_objects_yolo(img)
            
            # Store in index
            image_index[file_id] = {
                "name": image_file.filename,
                "embedding": embedding.cpu(),
                "objects": objects,
                "colors": colors,
                "folder": "Uploaded Images",
                "is_uploaded": True,
                "image_data": image_content  # Store image data for serving
            }
            
            # Create preview URL
            preview_url = f"/uploaded_image/{file_id}"
            
            uploaded_images.append({
                "file_id": file_id,
                "name": image_file.filename,
                "objects": objects,
                "colors": colors,
                "preview_url": preview_url
            })
            
            print(f"Uploaded and indexed: {image_file.filename} - Objects: {objects} - Colors: {colors}")
            
        except Exception as e:
            print(f"Failed to process uploaded image {image_file.filename}: {e}")
            continue
    
    return JSONResponse(content={
        "count": len(uploaded_images),
        "uploaded_images": uploaded_images,
        "message": f"Successfully uploaded and indexed {len(uploaded_images)} image(s)"
    })

@app.get("/uploaded_image/{file_id}")
async def get_uploaded_image(file_id: str):
    """Serve uploaded image"""
    if file_id not in image_index:
        raise HTTPException(status_code=404, detail="Image not found")
    
    image_data = image_index[file_id]
    if not image_data.get("is_uploaded"):
        raise HTTPException(status_code=404, detail="Image not found")
    
    image_content = image_data["image_data"]
    
    return StreamingResponse(
        io.BytesIO(image_content),
        media_type="image/jpeg",
        headers={"Content-Disposition": f"inline; filename={image_data['name']}"}
    )

@app.post("/analyze_storyboard")
async def analyze_storyboard(storyboard: UploadFile = File(...), guidelines: str = Form(""), feedback_images: str = Form(""), negative_feedback: str = Form("")):
    """Analyze storyboard image and find similar images with re-search capability"""
    global image_index
    
    if not image_index:
        raise HTTPException(status_code=400, detail="No images indexed. Please index your images first.")
    
    try:
        # Read and process storyboard image
        image_content = await storyboard.read()
        image_bytes = io.BytesIO(image_content)
        img = Image.open(image_bytes).convert("RGB")
        
        # Analyze storyboard
        storyboard_analysis = analyze_storyboard_image(img)
        
        # Find similar images
        similar_images = []
        for file_id, data in image_index.items():
            # Calculate semantic similarity
            semantic_sim = (storyboard_analysis['embedding'] @ data['embedding'].T).item()
            
            # Calculate object similarity
            storyboard_objects = set(storyboard_analysis['objects'])
            image_objects = set(data['objects'])
            if storyboard_objects or image_objects:
                object_sim = len(storyboard_objects & image_objects) / len(storyboard_objects | image_objects)
            else:
                object_sim = 1.0
            
            # Calculate color similarity
            color_sim = color_match_score(data['colors'], storyboard_analysis['colors'])
            
            # Combined similarity score
            similarity_score = calculate_combined_score(semantic_sim, object_sim, color_sim)
            
            similar_images.append({
                'file_id': file_id,
                'name': data['name'],
                'folder': data.get('folder', 'Root'),
                'similarity_score': similarity_score,
                'object_match': object_sim,
                'color_match': color_sim,
                'objects': data['objects'],
                'colors': data['colors']
            })
        
        # Sort by similarity score
        similar_images.sort(key=lambda x: x['similarity_score'], reverse=True)
        
        # Return top 10 similar images
        return JSONResponse(content={
            'analysis': {
                'objects': storyboard_analysis['objects'],
                'colors': storyboard_analysis['colors'],
                'suggested_rooms': storyboard_analysis['suggested_rooms']
            },
            'similar_images': similar_images[:10],
            'message': f"Found {len(similar_images)} similar images"
        })
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Storyboard analysis failed: {str(e)}")

@app.get("/")
def root():
    """API root endpoint - returns available endpoints"""
    return {
        "message": "Google Drive AI Search v3 with YOLOv8 running",
        "frontend": "Access the UI at http://localhost:4000/",
        "endpoints": {
            "auth": "/auth - Authenticate with Google Drive",
            "auth_status": "/auth/status - Check authentication status",
            "auth_disconnect": "/auth/disconnect - Disconnect from Google Drive",
            "index": "/index - Index all Drive images",
            "search": "/search - Search images with AI",
            "search_with_feedback": "/search_with_feedback - Advanced search with feedback",
            "parse": "/parse_requirements - Parse storyboard/PDF",
            "upload": "/upload_images - Upload and index images",
            "storyboard": "/analyze_storyboard - Analyze storyboard and find similar images",
            "stats": "/stats - Get indexing statistics",
            "image": "/image/{file_id} - Get image from Drive",
            "uploaded": "/uploaded_image/{file_id} - Get uploaded image",
            "export": "/export_pdf - Export images to PDF",
            "health": "/health - Health check",
            "test_supabase": "/test_supabase - Test Supabase connection",
            "setup_supabase": "/setup_supabase_table - Setup Supabase table",
            "clear_supabase": "/clear_supabase - Clear Supabase data"
        }
    }

@app.get("/debug_drive")
def debug_drive():
    """Debug endpoint to see what files the service account can access"""
    if not drive_service:
        return {"error": "Not authenticated"}
    
    try:
        # Test basic file listing
        results = drive_service.files().list(pageSize=10, fields="files(id,name,mimeType)").execute()
        files = results.get('files', [])
        
        # Test image-specific query
        image_results = drive_service.files().list(
            q="mimeType contains 'image/' and trashed=false", 
            pageSize=10, 
            fields="files(id,name,mimeType)"
        ).execute()
        image_files = image_results.get('files', [])
        
        return {
            "total_files_accessible": len(files),
            "sample_files": files[:3],
            "total_images_accessible": len(image_files),
            "sample_images": image_files[:3],
            "query_used": "mimeType contains 'image/' and trashed=false"
        }
    except Exception as e:
        return {"error": str(e)}

@app.get("/api")
def api_info():
    """Get API endpoints information"""
    return {
        "message": "Google Drive AI Search v3 with YOLOv8 running",
        "endpoints": {
            "auth": "/auth - Authenticate with Google Drive",
            "auth_status": "/auth/status - Check authentication status",
            "auth_disconnect": "/auth/disconnect - Disconnect from Google Drive",
            "index": "/index - Index all Drive images",
            "search": "/search - Search images with AI",
            "search_with_feedback": "/search_with_feedback - Advanced search with feedback",
            "parse": "/parse_requirements - Parse storyboard/PDF",
            "upload": "/upload_images - Upload and index images",
            "storyboard": "/analyze_storyboard - Analyze storyboard and find similar images",
            "stats": "/stats - Get indexing statistics",
            "image": "/image/{file_id} - Get image from Drive",
            "uploaded": "/uploaded_image/{file_id} - Get uploaded image",
            "export": "/export_pdf - Export images to PDF",
            "health": "/health - Health check",
            "test_supabase": "/test_supabase - Test Supabase connection",
            "setup_supabase": "/setup_supabase_table - Setup Supabase table",
            "clear_supabase": "/clear_supabase - Clear Supabase data"
        }
    }
